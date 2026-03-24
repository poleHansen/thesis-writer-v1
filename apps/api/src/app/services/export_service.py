from __future__ import annotations

from io import BytesIO
from pathlib import Path

import cairosvg
from pypdf import PdfWriter
from pptx import Presentation
from pptx.util import Inches


class PptxExportService:
    _SLIDE_WIDTH_INCHES = 13.333
    _SLIDE_HEIGHT_INCHES = 7.5

    def export_svg_pages_to_pptx(self, svg_paths: list[Path], target_path: str | Path) -> dict[str, object]:
        presentation = Presentation()
        presentation.slide_width = Inches(self._SLIDE_WIDTH_INCHES)
        presentation.slide_height = Inches(self._SLIDE_HEIGHT_INCHES)

        blank_layout = presentation.slide_layouts[6]
        rendered_files: list[str] = []

        for index, svg_path in enumerate(svg_paths):
            if index == 0 and len(presentation.slides) == 1 and not presentation.slides[0].shapes:
                slide = presentation.slides[0]
            else:
                slide = presentation.slides.add_slide(blank_layout)

            png_bytes = cairosvg.svg2png(url=str(svg_path))
            picture_stream = BytesIO(png_bytes)
            slide.shapes.add_picture(
                picture_stream,
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
            rendered_files.append(svg_path.name)

        target = Path(target_path)
        target.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(target)

        return {
            "file_count": len(svg_paths),
            "rendered_files": rendered_files,
            "renderer": "cairosvg+python-pptx",
        }


class PdfExportService:
    def export_svg_pages_to_pdf(self, svg_paths: list[Path], target_path: str | Path) -> dict[str, object]:
        writer = PdfWriter()
        rendered_files: list[str] = []

        for svg_path in svg_paths:
            pdf_bytes = cairosvg.svg2pdf(url=str(svg_path))
            page_reader = PdfWriter(clone_from=BytesIO(pdf_bytes))
            for page in page_reader.pages:
                writer.add_page(page)
            rendered_files.append(svg_path.name)

        target = Path(target_path)
        target.parent.mkdir(parents=True, exist_ok=True)
        with target.open("wb") as output_stream:
            writer.write(output_stream)

        return {
            "file_count": len(svg_paths),
            "rendered_files": rendered_files,
            "renderer": "cairosvg+pypdf",
        }