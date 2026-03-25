from __future__ import annotations

import base64
import json
import os
import sys
from pathlib import Path

from fastapi.testclient import TestClient
from pptx import Presentation
from pypdf import PdfReader


def test_generation_writes_artifact_snapshots(tmp_path: Path) -> None:
    database_path = tmp_path / "artifact-test.db"
    storage_root = Path("storage")

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Artifact Persistence Demo",
        "description": "Verify generated artifacts are written to disk.",
        "source_mode": "file",
        "tags": ["test"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "demo-checksum",
        "extracted_summary": "A compact source summary for artifact tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(
            f"/projects/{project_id}/brief:generate",
            json={"force_regenerate": True},
        )
        assert brief_response.status_code == 200
        brief_id = brief_response.json()["brief"]["id"]

        outline_response = client.post(
            f"/projects/{project_id}/outline:generate",
            json={"brief_id": brief_id},
        )
        assert outline_response.status_code == 200
        outline_id = outline_response.json()["outline"]["id"]

        slide_plan_response = client.post(
            f"/projects/{project_id}/slide-plan:generate",
            json={"outline_id": outline_id},
        )
        assert slide_plan_response.status_code == 200
        slide_plan_id = slide_plan_response.json()["slide_plan"]["id"]

        render_response = client.post(
            f"/projects/{project_id}/artifact:generate",
            json={"slide_plan_id": slide_plan_id},
        )
        assert render_response.status_code == 200
        artifact_id = render_response.json()["artifact"]["id"]

    artifact_dir = storage_root / "projects" / project_id / "artifacts"
    brief_path = artifact_dir / "brief.json"
    outline_path = artifact_dir / "outline.json"
    slide_plan_path = artifact_dir / "slide-plan.json"
    slide_artifact_path = artifact_dir / "slide-artifact.json"
    render_root = storage_root / "projects" / project_id / "render" / artifact_id
    design_spec_path = render_root / "design-spec.json"

    assert brief_path.exists()
    assert outline_path.exists()
    assert slide_plan_path.exists()
    assert slide_artifact_path.exists()
    assert (render_root / "svg_output").exists()
    assert (render_root / "svg_final").exists()
    assert design_spec_path.exists()
    assert (render_root / "render-log.json").exists()
    assert any((render_root / "svg_final").glob("slide-*.svg"))

    brief_payload = json.loads(brief_path.read_text(encoding="utf-8"))
    outline_payload = json.loads(outline_path.read_text(encoding="utf-8"))
    slide_plan_payload = json.loads(slide_plan_path.read_text(encoding="utf-8"))
    slide_artifact_payload = json.loads(slide_artifact_path.read_text(encoding="utf-8"))
    design_spec_payload = json.loads(design_spec_path.read_text(encoding="utf-8"))
    render_log_payload = json.loads((render_root / "render-log.json").read_text(encoding="utf-8"))
    first_svg = next((render_root / "svg_final").glob("slide-*.svg"))
    first_svg_content = first_svg.read_text(encoding="utf-8")

    assert brief_payload["project_id"] == project_id
    assert outline_payload["project_id"] == project_id
    assert slide_plan_payload["project_id"] == project_id
    assert slide_artifact_payload["project_id"] == project_id
    assert slide_artifact_payload["slide_plan_id"] == slide_plan_id
    assert slide_artifact_payload["metadata"]["design_spec_path"].endswith("design-spec.json")
    assert slide_artifact_payload["render_status"] in {"succeeded", "partial"}
    assert slide_artifact_payload["metadata"]["generated_svg_files"]
    assert slide_artifact_payload["metadata"]["validation_summary"]["checked_file_count"] >= 1
    assert slide_artifact_payload["metadata"]["validation_summary"]["invalid_file_count"] == 0
    assert slide_artifact_payload["metadata"]["finalization_summary"]["page_count"] >= 1
    assert slide_artifact_payload["metadata"]["finalization_summary"]["width_height_alignment_count"] >= 0
    assert design_spec_payload["project_id"] == project_id
    assert design_spec_payload["page_count"] >= 1
    assert design_spec_payload["template"]["template_id"]
    assert slide_plan_payload["page_count"] >= 1
    assert render_log_payload["validation_results"]
    assert all(item["is_valid"] for item in render_log_payload["validation_results"])
    assert render_log_payload["finalization_summary"]["page_count"] >= 1
    assert render_log_payload["finalization_summary"]["width_height_alignment_count"] >= 0
    assert all(item["finalizer_steps"] for item in render_log_payload["validation_results"])
    assert all(
        "ensure_canvas_dimensions" in item["finalizer_steps"] or 'width="1280"' in first_svg_content
        for item in render_log_payload["validation_results"]
    )
    assert first_svg_content.startswith("<?xml version=\"1.0\" encoding=\"UTF-8\"?>")
    assert 'width="1280"' in first_svg_content
    assert 'height="720"' in first_svg_content


def test_export_writes_project_export_artifact(tmp_path: Path) -> None:
    database_path = tmp_path / "export-test.db"
    storage_root = Path("storage")

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Export Demo",
        "description": "Verify export jobs persist and write output.",
        "source_mode": "file",
        "tags": ["test"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "demo-checksum",
        "extracted_summary": "A compact source summary for export tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief_id = brief_response.json()["brief"]["id"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief_id})
        assert outline_response.status_code == 200
        outline_id = outline_response.json()["outline"]["id"]

        slide_plan_response = client.post(f"/projects/{project_id}/slide-plan:generate", json={"outline_id": outline_id})
        assert slide_plan_response.status_code == 200
        slide_plan_id = slide_plan_response.json()["slide_plan"]["id"]

        render_response = client.post(f"/projects/{project_id}/artifact:generate", json={"slide_plan_id": slide_plan_id})
        assert render_response.status_code == 200
        artifact_id = render_response.json()["artifact"]["id"]

        export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact_id, "export_format": "pptx"},
        )
        assert export_response.status_code == 200
        export_payload = export_response.json()

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

    export_job = export_payload["export_job"]
    export_path = Path(export_job["export_path"])
    archive_manifest_path = Path(export_job["metadata"]["archive_manifest_path"])
    export_log_path = Path(export_job["metadata"]["export_log_path"])
    assert export_path.exists()
    assert export_path.suffix == ".pptx"
    assert archive_manifest_path.exists()
    assert export_log_path.exists()
    presentation = Presentation(export_path)
    archive_manifest = json.loads(archive_manifest_path.read_text(encoding="utf-8"))
    export_log = json.loads(export_log_path.read_text(encoding="utf-8"))
    assert export_job["status"] == "succeeded"
    assert export_job["run_id"].startswith("run-")
    assert export_job["metadata"]["export_kind"] == "pptx_from_svg_pages"
    assert export_job["metadata"]["file_count"] >= 1
    assert export_job["metadata"]["renderer"] == "cairosvg+python-pptx"
    assert export_job["metadata"]["run_id"] == export_job["run_id"]
    assert export_job["run_id"] in export_job["export_path"]
    assert len(export_job["metadata"]["rendered_files"]) == len(presentation.slides)
    assert len(presentation.slides) == export_job["metadata"]["file_count"]
    assert archive_manifest["project_id"] == project_id
    assert archive_manifest["run_id"] == export_job["run_id"]
    assert archive_manifest["export_path"] == export_job["export_path"]
    assert archive_manifest["artifacts"]["export_output"] == export_job["export_path"]
    assert archive_manifest["artifacts"]["slide_artifact"].endswith("slide-artifact.json")
    assert export_log["run_id"] == export_job["run_id"]
    assert export_log["archive_manifest_path"] == export_job["metadata"]["archive_manifest_path"]
    assert export_log["status"] == "succeeded"
    assert export_payload["task_run"]["task_type"] == "export"
    assert export_payload["task_run"]["task_status"] == "succeeded"
    assert export_payload["task_run"]["result"]["run_id"] == export_job["run_id"]
    assert export_payload["task_run"]["result"]["archive_manifest_path"] == export_job["metadata"]["archive_manifest_path"]
    assert detail_payload["project"]["status"] == "exported"
    assert detail_payload["latest_export"]["id"] == export_job["id"]
    assert detail_payload["latest_export"]["run_id"] == export_job["run_id"]
    assert detail_payload["latest_export"]["export_path"] == export_job["export_path"]


def test_export_failure_is_persisted_for_traceability(tmp_path: Path) -> None:
    database_path = tmp_path / "export-failure-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Export Failure Demo",
        "description": "Verify failed export attempts are persisted.",
        "source_mode": "file",
        "tags": ["test"],
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": "missing-artifact", "export_format": "pptx"},
        )
        assert export_response.status_code == 400
        assert export_response.json()["detail"] == "Project has no rendered artifact to export"

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        status_response = client.get(f"/projects/{project_id}/status")
        assert status_response.status_code == 200
        status_payload = status_response.json()

    assert detail_payload["project"]["status"] == "export_failed"
    assert detail_payload["latest_export"] is None
    assert status_payload["current_task"]["task_type"] == "export"
    assert status_payload["current_task"]["task_status"] == "failed"
    assert status_payload["current_task"]["error_message"] == "Project has no rendered artifact to export"


def test_pdf_preview_export_writes_preview_path(tmp_path: Path) -> None:
    database_path = tmp_path / "export-pdf-test.db"
    storage_root = Path("storage")

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Export PDF Demo",
        "description": "Verify PDF preview exports persist and write output.",
        "source_mode": "file",
        "tags": ["test"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "demo-checksum",
        "extracted_summary": "A compact source summary for PDF export tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief_id = brief_response.json()["brief"]["id"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief_id})
        assert outline_response.status_code == 200
        outline_id = outline_response.json()["outline"]["id"]

        slide_plan_response = client.post(f"/projects/{project_id}/slide-plan:generate", json={"outline_id": outline_id})
        assert slide_plan_response.status_code == 200
        slide_plan_id = slide_plan_response.json()["slide_plan"]["id"]

        render_response = client.post(f"/projects/{project_id}/artifact:generate", json={"slide_plan_id": slide_plan_id})
        assert render_response.status_code == 200
        artifact_id = render_response.json()["artifact"]["id"]

        export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact_id, "export_format": "pdf"},
        )
        assert export_response.status_code == 200
        export_payload = export_response.json()

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

    export_job = export_payload["export_job"]
    preview_pdf_path = Path(export_job["preview_pdf_path"])
    archive_manifest_path = Path(export_job["metadata"]["archive_manifest_path"])
    assert storage_root.exists()
    assert preview_pdf_path.exists()
    assert archive_manifest_path.exists()
    assert preview_pdf_path.suffix == ".pdf"
    reader = PdfReader(str(preview_pdf_path))
    archive_manifest = json.loads(archive_manifest_path.read_text(encoding="utf-8"))
    assert len(reader.pages) == export_job["metadata"]["file_count"]
    assert export_job["status"] == "succeeded"
    assert export_job["run_id"].startswith("run-")
    assert export_job["metadata"]["export_kind"] == "pdf_preview_from_svg_pages"
    assert export_job["metadata"]["renderer"] == "cairosvg+pypdf"
    assert export_job["export_path"] == export_job["preview_pdf_path"]
    assert archive_manifest["run_id"] == export_job["run_id"]
    assert archive_manifest["preview_pdf_path"] == export_job["preview_pdf_path"]
    assert archive_manifest["export_format"] == "pdf"
    assert detail_payload["project"]["status"] == "exported"
    assert detail_payload["latest_export"]["id"] == export_job["id"]
    assert detail_payload["latest_export"]["preview_pdf_path"] == export_job["preview_pdf_path"]


def test_repeated_exports_create_distinct_run_ids_and_paths(tmp_path: Path) -> None:
    database_path = tmp_path / "export-run-id-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Export Run ID Demo",
        "description": "Verify repeated exports produce distinct run ids and output paths.",
        "source_mode": "file",
        "tags": ["test"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "demo-checksum",
        "extracted_summary": "A compact source summary for run id tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief_id = brief_response.json()["brief"]["id"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief_id})
        assert outline_response.status_code == 200
        outline_id = outline_response.json()["outline"]["id"]

        slide_plan_response = client.post(f"/projects/{project_id}/slide-plan:generate", json={"outline_id": outline_id})
        assert slide_plan_response.status_code == 200
        slide_plan_id = slide_plan_response.json()["slide_plan"]["id"]

        render_response = client.post(f"/projects/{project_id}/artifact:generate", json={"slide_plan_id": slide_plan_id})
        assert render_response.status_code == 200
        artifact_id = render_response.json()["artifact"]["id"]

        first_export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact_id, "export_format": "pptx"},
        )
        assert first_export_response.status_code == 200
        first_export_job = first_export_response.json()["export_job"]

        second_export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact_id, "export_format": "pptx"},
        )
        assert second_export_response.status_code == 200
        second_export_job = second_export_response.json()["export_job"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

    assert first_export_job["run_id"] != second_export_job["run_id"]
    assert first_export_job["export_path"] != second_export_job["export_path"]
    assert Path(first_export_job["export_path"]).exists()
    assert Path(second_export_job["export_path"]).exists()
    assert first_export_job["run_id"] in first_export_job["export_path"]
    assert second_export_job["run_id"] in second_export_job["export_path"]
    assert detail_payload["latest_export"]["id"] == second_export_job["id"]
    assert detail_payload["latest_export"]["run_id"] == second_export_job["run_id"]


def test_list_project_exports_returns_recent_runs_in_desc_order(tmp_path: Path) -> None:
    database_path = tmp_path / "export-history-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Export History Demo",
        "description": "Verify recent export runs are listed.",
        "source_mode": "file",
        "tags": ["test"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "demo-checksum",
        "extracted_summary": "A compact source summary for export history tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief_id = brief_response.json()["brief"]["id"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief_id})
        assert outline_response.status_code == 200
        outline_id = outline_response.json()["outline"]["id"]

        slide_plan_response = client.post(f"/projects/{project_id}/slide-plan:generate", json={"outline_id": outline_id})
        assert slide_plan_response.status_code == 200
        slide_plan_id = slide_plan_response.json()["slide_plan"]["id"]

        render_response = client.post(f"/projects/{project_id}/artifact:generate", json={"slide_plan_id": slide_plan_id})
        assert render_response.status_code == 200
        artifact_id = render_response.json()["artifact"]["id"]

        run_ids = []
        for export_format in ["pptx", "pdf", "pptx"]:
          export_response = client.post(
              f"/projects/{project_id}/export",
              json={"artifact_id": artifact_id, "export_format": export_format},
          )
          assert export_response.status_code == 200
          run_ids.append(export_response.json()["export_job"]["run_id"])

        history_response = client.get(f"/projects/{project_id}/exports?limit=2")
        assert history_response.status_code == 200
        history_payload = history_response.json()

    assert len(history_payload["exports"]) == 2
    assert history_payload["exports"][0]["run_id"] == run_ids[2]
    assert history_payload["exports"][1]["run_id"] == run_ids[1]


def test_list_project_exports_honors_positive_limit(tmp_path: Path) -> None:
    database_path = tmp_path / "export-history-positive-limit-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Export History Positive Limit",
        "description": "Verify export history only returns the requested number of latest runs.",
        "source_mode": "file",
        "tags": ["test"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "demo-checksum-positive-limit",
        "extracted_summary": "A compact source summary for positive export history limit tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief_id = brief_response.json()["brief"]["id"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief_id})
        assert outline_response.status_code == 200
        outline_id = outline_response.json()["outline"]["id"]

        slide_plan_response = client.post(f"/projects/{project_id}/slide-plan:generate", json={"outline_id": outline_id})
        assert slide_plan_response.status_code == 200
        slide_plan_id = slide_plan_response.json()["slide_plan"]["id"]

        render_response = client.post(f"/projects/{project_id}/artifact:generate", json={"slide_plan_id": slide_plan_id})
        assert render_response.status_code == 200
        artifact_id = render_response.json()["artifact"]["id"]

        run_ids = []
        for export_format in ["pptx", "pdf", "pptx"]:
            export_response = client.post(
                f"/projects/{project_id}/export",
                json={"artifact_id": artifact_id, "export_format": export_format},
            )
            assert export_response.status_code == 200
            run_ids.append(export_response.json()["export_job"]["run_id"])

        history_response = client.get(f"/projects/{project_id}/exports?limit=1")
        assert history_response.status_code == 200
        history_payload = history_response.json()

    assert history_payload["exports"]
    assert len(history_payload["exports"]) == 1
    assert history_payload["exports"][0]["run_id"] == run_ids[2]


def test_list_projects_dashboard_keeps_early_parse_states_isolated(tmp_path: Path) -> None:
    database_path = tmp_path / "dashboard-early-parse-states.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()

    parsed_upload = {
        "file_name": "parsed.md",
        "file_type": "markdown",
        "mime_type": "text/markdown",
        "content_base64": base64.b64encode(b"# Parsed Demo\n\nThis file produces a clean parsed bundle.").decode("ascii"),
        "metadata": {},
    }
    analyzed_upload = {
        "file_name": "analyzed.md",
        "file_type": "markdown",
        "mime_type": "text/markdown",
        "content_base64": base64.b64encode(b"   \n\n   ").decode("ascii"),
        "metadata": {},
    }
    failed_file_payload = {
        "file_name": "parse-failed.xyz",
        "file_type": "auto_detected",
        "storage_path": "storage/uploads/parse-failed.xyz",
        "mime_type": "application/octet-stream",
        "size_bytes": 64,
        "checksum": "parse-failed-dashboard-checksum",
        "extracted_summary": "A placeholder file record that cannot be parsed by the current parser.",
        "metadata": {},
    }

    with TestClient(app) as client:
        parsed_project = client.post(
            "/projects",
            json={
                "name": "Parsed Dashboard Demo",
                "description": "Verify dashboard keeps parsed projects isolated.",
                "source_mode": "file",
                "tags": ["phase7", "dashboard", "parsed"],
            },
        )
        assert parsed_project.status_code == 201
        parsed_project_id = parsed_project.json()["project"]["id"]
        parsed_file = client.post(f"/projects/{parsed_project_id}/files:upload", json=parsed_upload)
        assert parsed_file.status_code == 201
        parsed_parse = client.post(
            f"/projects/{parsed_project_id}/files:parse",
            json={"file_ids": [parsed_file.json()["file"]["id"]], "rebuild_bundle": True},
        )
        assert parsed_parse.status_code == 200
        parsed_bundle = parsed_parse.json()["source_bundle"]

        analyzed_project = client.post(
            "/projects",
            json={
                "name": "Analyzed Dashboard Demo",
                "description": "Verify dashboard keeps analyzed projects isolated.",
                "source_mode": "file",
                "tags": ["phase7", "dashboard", "analyzed"],
            },
        )
        assert analyzed_project.status_code == 201
        analyzed_project_id = analyzed_project.json()["project"]["id"]
        analyzed_file = client.post(f"/projects/{analyzed_project_id}/files:upload", json=analyzed_upload)
        assert analyzed_file.status_code == 201
        analyzed_parse = client.post(
            f"/projects/{analyzed_project_id}/files:parse",
            json={"file_ids": [analyzed_file.json()["file"]["id"]], "rebuild_bundle": True},
        )
        assert analyzed_parse.status_code == 200
        analyzed_bundle = analyzed_parse.json()["source_bundle"]

        failed_project = client.post(
            "/projects",
            json={
                "name": "Parse Failed Dashboard Demo",
                "description": "Verify dashboard keeps parse_failed projects isolated.",
                "source_mode": "file",
                "tags": ["phase7", "dashboard", "parse_failed"],
            },
        )
        assert failed_project.status_code == 201
        failed_project_id = failed_project.json()["project"]["id"]
        failed_file = client.post(f"/projects/{failed_project_id}/files", json=failed_file_payload)
        assert failed_file.status_code == 201
        failed_parse = client.post(
            f"/projects/{failed_project_id}/files:parse",
            json={"file_ids": [failed_file.json()["file"]["id"]], "rebuild_bundle": True},
        )
        assert failed_parse.status_code == 200

        dashboard_response = client.get("/projects")
        assert dashboard_response.status_code == 200
        dashboard_payload = dashboard_response.json()

    items_by_id = {item["project"]["id"]: item for item in dashboard_payload["projects"]}

    parsed_item = items_by_id[parsed_project_id]
    assert parsed_item["project"]["status"] == "parsed"
    assert parsed_item["file_count"] == 1
    assert parsed_item["parsed_file_count"] == 1
    assert parsed_item["failed_file_count"] == 0
    assert parsed_item["latest_brief"] is None
    assert parsed_item["latest_outline"] is None
    assert parsed_item["latest_slide_plan"] is None
    assert parsed_item["latest_artifact"] is None
    assert parsed_item["latest_export"] is None
    assert parsed_item["current_task"]["task_type"] == "parse"
    assert parsed_item["current_task"]["task_status"] == "succeeded"
    assert parsed_item["current_task"]["result"]["source_bundle_id"] == parsed_bundle["id"]

    analyzed_item = items_by_id[analyzed_project_id]
    assert analyzed_item["project"]["status"] == "analyzed"
    assert analyzed_item["file_count"] == 1
    assert analyzed_item["parsed_file_count"] == 1
    assert analyzed_item["failed_file_count"] == 0
    assert analyzed_item["latest_brief"] is None
    assert analyzed_item["latest_outline"] is None
    assert analyzed_item["latest_slide_plan"] is None
    assert analyzed_item["latest_artifact"] is None
    assert analyzed_item["latest_export"] is None
    assert analyzed_item["current_task"]["task_type"] == "parse"
    assert analyzed_item["current_task"]["task_status"] == "succeeded"
    assert analyzed_item["current_task"]["result"]["source_bundle_id"] == analyzed_bundle["id"]

    failed_item = items_by_id[failed_project_id]
    assert failed_item["project"]["status"] == "parse_failed"
    assert failed_item["file_count"] == 1
    assert failed_item["parsed_file_count"] == 0
    assert failed_item["failed_file_count"] == 1
    assert failed_item["latest_brief"] is None
    assert failed_item["latest_outline"] is None
    assert failed_item["latest_slide_plan"] is None
    assert failed_item["latest_artifact"] is None
    assert failed_item["latest_export"] is None
    assert failed_item["current_task"]["task_type"] == "parse"
    assert failed_item["current_task"]["task_status"] == "failed"
    assert failed_item["current_task"]["result"]["source_bundle_id"] is None


def test_list_projects_dashboard_keeps_briefing_outlined_and_planned_states_isolated(tmp_path: Path) -> None:
    database_path = tmp_path / "dashboard-briefing-outlined-planned.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 192,
        "checksum": "dashboard-briefing-outlined-planned-checksum",
        "extracted_summary": "A compact source summary for dashboard planning-state isolation tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        briefing_response = client.post(
            "/projects",
            json={
                "name": "Dashboard Briefing Project",
                "description": "Verify dashboard keeps briefing projects isolated.",
                "source_mode": "file",
                "tags": ["phase7", "dashboard", "briefing"],
            },
        )
        assert briefing_response.status_code == 201
        briefing_project_id = briefing_response.json()["project"]["id"]

        briefing_register_response = client.post(f"/projects/{briefing_project_id}/files", json=file_payload)
        assert briefing_register_response.status_code == 201

        briefing_brief_response = client.post(
            f"/projects/{briefing_project_id}/brief:generate",
            json={"force_regenerate": True},
        )
        assert briefing_brief_response.status_code == 200
        briefing_brief = briefing_brief_response.json()["brief"]

        outlined_response = client.post(
            "/projects",
            json={
                "name": "Dashboard Outlined Project",
                "description": "Verify dashboard keeps outlined projects isolated.",
                "source_mode": "file",
                "tags": ["phase7", "dashboard", "outlined"],
            },
        )
        assert outlined_response.status_code == 201
        outlined_project_id = outlined_response.json()["project"]["id"]

        outlined_register_response = client.post(f"/projects/{outlined_project_id}/files", json=file_payload)
        assert outlined_register_response.status_code == 201

        outlined_brief_response = client.post(
            f"/projects/{outlined_project_id}/brief:generate",
            json={"force_regenerate": True},
        )
        assert outlined_brief_response.status_code == 200
        outlined_brief = outlined_brief_response.json()["brief"]

        outlined_outline_response = client.post(
            f"/projects/{outlined_project_id}/outline:generate",
            json={"brief_id": outlined_brief["id"]},
        )
        assert outlined_outline_response.status_code == 200
        outlined_outline = outlined_outline_response.json()["outline"]

        planned_response = client.post(
            "/projects",
            json={
                "name": "Dashboard Planned Project",
                "description": "Verify dashboard keeps planned projects isolated.",
                "source_mode": "file",
                "tags": ["phase7", "dashboard", "planned"],
            },
        )
        assert planned_response.status_code == 201
        planned_project_id = planned_response.json()["project"]["id"]

        planned_register_response = client.post(f"/projects/{planned_project_id}/files", json=file_payload)
        assert planned_register_response.status_code == 201

        planned_brief_response = client.post(
            f"/projects/{planned_project_id}/brief:generate",
            json={"force_regenerate": True},
        )
        assert planned_brief_response.status_code == 200
        planned_brief = planned_brief_response.json()["brief"]

        planned_outline_response = client.post(
            f"/projects/{planned_project_id}/outline:generate",
            json={"brief_id": planned_brief["id"]},
        )
        assert planned_outline_response.status_code == 200
        planned_outline = planned_outline_response.json()["outline"]

        planned_slide_plan_response = client.post(
            f"/projects/{planned_project_id}/slide-plan:generate",
            json={"outline_id": planned_outline["id"]},
        )
        assert planned_slide_plan_response.status_code == 200
        planned_slide_plan = planned_slide_plan_response.json()["slide_plan"]

        dashboard_response = client.get("/projects")
        assert dashboard_response.status_code == 200
        dashboard_payload = dashboard_response.json()

    assert len(dashboard_payload["projects"]) == 3

    newest_dashboard_item = dashboard_payload["projects"][0]
    middle_dashboard_item = dashboard_payload["projects"][1]
    oldest_dashboard_item = dashboard_payload["projects"][2]

    assert newest_dashboard_item["project"]["id"] == planned_project_id
    assert newest_dashboard_item["project"]["status"] == "planned"
    assert newest_dashboard_item["latest_brief"]["id"] == planned_brief["id"]
    assert newest_dashboard_item["latest_outline"]["id"] == planned_outline["id"]
    assert newest_dashboard_item["latest_slide_plan"]["id"] == planned_slide_plan["id"]
    assert newest_dashboard_item["latest_artifact"] is None
    assert newest_dashboard_item["latest_export"] is None
    assert newest_dashboard_item["current_task"]["task_type"] == "generate_slide_plan"
    assert newest_dashboard_item["current_task"]["task_status"] == "succeeded"
    assert newest_dashboard_item["current_task"]["project_id"] == planned_project_id

    assert middle_dashboard_item["project"]["id"] == outlined_project_id
    assert middle_dashboard_item["project"]["status"] == "outlined"
    assert middle_dashboard_item["latest_brief"]["id"] == outlined_brief["id"]
    assert middle_dashboard_item["latest_outline"]["id"] == outlined_outline["id"]
    assert middle_dashboard_item["latest_slide_plan"] is None
    assert middle_dashboard_item["latest_artifact"] is None
    assert middle_dashboard_item["latest_export"] is None
    assert middle_dashboard_item["current_task"]["task_type"] == "generate_outline"
    assert middle_dashboard_item["current_task"]["task_status"] == "succeeded"
    assert middle_dashboard_item["current_task"]["project_id"] == outlined_project_id

    assert oldest_dashboard_item["project"]["id"] == briefing_project_id
    assert oldest_dashboard_item["project"]["status"] == "briefing"
    assert oldest_dashboard_item["latest_brief"]["id"] == briefing_brief["id"]
    assert oldest_dashboard_item["latest_outline"] is None
    assert oldest_dashboard_item["latest_slide_plan"] is None
    assert oldest_dashboard_item["latest_artifact"] is None
    assert oldest_dashboard_item["latest_export"] is None
    assert oldest_dashboard_item["current_task"]["task_type"] == "generate_brief"
    assert oldest_dashboard_item["current_task"]["task_status"] == "succeeded"
    assert oldest_dashboard_item["current_task"]["project_id"] == briefing_project_id


def test_list_project_exports_returns_empty_list_when_project_has_no_exports(tmp_path: Path) -> None:
    database_path = tmp_path / "export-history-empty-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Export History Empty",
        "description": "Verify export history stays empty before any export run.",
        "source_mode": "file",
        "tags": ["test"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "demo-checksum-empty-history",
        "extracted_summary": "A compact source summary for empty export history tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief_id = brief_response.json()["brief"]["id"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief_id})
        assert outline_response.status_code == 200
        outline_id = outline_response.json()["outline"]["id"]

        slide_plan_response = client.post(f"/projects/{project_id}/slide-plan:generate", json={"outline_id": outline_id})
        assert slide_plan_response.status_code == 200
        slide_plan_id = slide_plan_response.json()["slide_plan"]["id"]

        render_response = client.post(f"/projects/{project_id}/artifact:generate", json={"slide_plan_id": slide_plan_id})
        assert render_response.status_code == 200

        history_response = client.get(f"/projects/{project_id}/exports?limit=5")
        assert history_response.status_code == 200
        history_payload = history_response.json()

    assert history_payload == {"exports": []}

def test_list_project_exports_rejects_non_positive_limit(tmp_path: Path) -> None:
    database_path = tmp_path / "export-history-limit-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    with TestClient(app) as client:
        create_response = client.post(
            "/projects",
            json={
                "name": "Export History Limit",
                "description": "Verify export history rejects non-positive limit.",
                "source_mode": "file",
                "tags": ["test"],
            },
        )
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        zero_limit_response = client.get(f"/projects/{project_id}/exports?limit=0")
        negative_limit_response = client.get(f"/projects/{project_id}/exports?limit=-1")

    assert zero_limit_response.status_code == 422
    assert negative_limit_response.status_code == 422


def test_get_project_export_returns_selected_run_and_rejects_foreign_project(tmp_path: Path) -> None:
    database_path = tmp_path / "export-detail-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    def build_project(client: TestClient, name: str) -> tuple[str, str]:
        project_response = client.post(
            "/projects",
            json={
                "name": name,
                "description": "Export detail lookup test.",
                "source_mode": "file",
                "tags": ["test"],
            },
        )
        assert project_response.status_code == 201
        project_id = project_response.json()["project"]["id"]

        file_response = client.post(
            f"/projects/{project_id}/files",
            json={
                "file_name": "demo.md",
                "file_type": "markdown",
                "storage_path": "storage/uploads/demo.md",
                "mime_type": "text/markdown",
                "size_bytes": 128,
                "checksum": f"checksum-{name}",
                "extracted_summary": "A compact source summary for export detail tests.",
                "metadata": {},
            },
        )
        assert file_response.status_code == 201

        brief_id = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True}).json()["brief"]["id"]
        outline_id = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief_id}).json()["outline"]["id"]
        slide_plan_id = client.post(f"/projects/{project_id}/slide-plan:generate", json={"outline_id": outline_id}).json()["slide_plan"]["id"]
        artifact_id = client.post(f"/projects/{project_id}/artifact:generate", json={"slide_plan_id": slide_plan_id}).json()["artifact"]["id"]

        export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact_id, "export_format": "pptx"},
        )
        assert export_response.status_code == 200
        export_id = export_response.json()["export_job"]["id"]
        return project_id, export_id

    with TestClient(app) as client:
        first_project_id, first_export_id = build_project(client, "Export Detail One")
        second_project_id, _ = build_project(client, "Export Detail Two")

        detail_response = client.get(f"/projects/{first_project_id}/exports/{first_export_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        missing_export_response = client.get(f"/projects/{first_project_id}/exports/missing-export")
        assert missing_export_response.status_code == 404
        assert missing_export_response.json() == {"detail": "Export job not found"}

        foreign_response = client.get(f"/projects/{second_project_id}/exports/{first_export_id}")
        assert foreign_response.status_code == 404

    assert detail_payload["export_job"]["id"] == first_export_id
    assert detail_payload["export_job"]["project_id"] == first_project_id


def test_get_project_export_treats_malformed_but_routable_id_as_not_found(tmp_path: Path) -> None:
    database_path = tmp_path / "export-detail-malformed-id-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    with TestClient(app) as client:
        create_response = client.post(
            "/projects",
            json={
                "name": "Export Detail Malformed",
                "description": "Verify malformed-but-routable export ids stay in not-found semantics.",
                "source_mode": "file",
                "tags": ["test"],
            },
        )
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        response = client.get(f"/projects/{project_id}/exports/%20%20not-a-real-export%20%20")

    assert response.status_code == 404
    assert response.json() == {"detail": "Export job not found"}


def test_list_projects_dashboard_returns_aggregated_project_summary(tmp_path: Path) -> None:
    database_path = tmp_path / "dashboard-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Dashboard Demo",
        "description": "Verify dashboard aggregation returns project summary.",
        "source_mode": "file",
        "tags": ["phase6", "dashboard"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "demo-checksum",
        "extracted_summary": "A compact source summary for dashboard tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project = create_response.json()["project"]
        project_id = project["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief = brief_response.json()["brief"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief["id"]})
        assert outline_response.status_code == 200
        outline = outline_response.json()["outline"]

        slide_plan_response = client.post(f"/projects/{project_id}/slide-plan:generate", json={"outline_id": outline["id"]})
        assert slide_plan_response.status_code == 200
        slide_plan = slide_plan_response.json()["slide_plan"]

        artifact_response = client.post(f"/projects/{project_id}/artifact:generate", json={"slide_plan_id": slide_plan["id"]})
        assert artifact_response.status_code == 200
        artifact = artifact_response.json()["artifact"]

        export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact["id"], "export_format": "pdf"},
        )
        assert export_response.status_code == 200
        export_job = export_response.json()["export_job"]

        dashboard_response = client.get("/projects")
        assert dashboard_response.status_code == 200
        dashboard_payload = dashboard_response.json()

    assert len(dashboard_payload["projects"]) == 1
    dashboard_item = dashboard_payload["projects"][0]
    assert dashboard_item["project"]["id"] == project_id
    assert dashboard_item["project"]["name"] == project_payload["name"]
    assert dashboard_item["project"]["status"] == "exported"
    assert dashboard_item["file_count"] == 1
    assert dashboard_item["parsed_file_count"] == 0
    assert dashboard_item["failed_file_count"] == 0
    assert dashboard_item["latest_brief"]["id"] == brief["id"]
    assert dashboard_item["latest_outline"]["id"] == outline["id"]
    assert dashboard_item["latest_slide_plan"]["id"] == slide_plan["id"]
    assert dashboard_item["latest_artifact"]["id"] == artifact["id"]
    assert dashboard_item["latest_export"]["id"] == export_job["id"]
    assert dashboard_item["latest_export"]["preview_pdf_path"] == export_job["preview_pdf_path"]
    assert dashboard_item["current_task"]["task_type"] == "export"
    assert dashboard_item["current_task"]["task_status"] == "succeeded"


def test_list_projects_dashboard_keeps_single_created_project_empty(tmp_path: Path) -> None:
    database_path = tmp_path / "dashboard-single-created-project.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Dashboard Single Created Project",
        "description": "Verify dashboard keeps a newly created project empty and aligned.",
        "source_mode": "chat",
        "tags": ["phase7", "dashboard", "created"],
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        dashboard_response = client.get("/projects")
        assert dashboard_response.status_code == 200
        dashboard_payload = dashboard_response.json()

    assert len(dashboard_payload["projects"]) == 1
    dashboard_item = dashboard_payload["projects"][0]
    assert dashboard_item["project"]["id"] == project_id
    assert dashboard_item["project"]["status"] == "created"
    assert dashboard_item["file_count"] == 0
    assert dashboard_item["parsed_file_count"] == 0
    assert dashboard_item["failed_file_count"] == 0
    assert dashboard_item["latest_brief"] is None
    assert dashboard_item["latest_outline"] is None
    assert dashboard_item["latest_slide_plan"] is None
    assert dashboard_item["latest_artifact"] is None
    assert dashboard_item["latest_export"] is None
    assert dashboard_item["current_task"]["task_type"] == "ingest"
    assert dashboard_item["current_task"]["task_status"] == "succeeded"
    assert dashboard_item["current_task"]["project_id"] == project_id


def test_list_projects_dashboard_keeps_latest_export_on_most_recent_run(tmp_path: Path) -> None:
    database_path = tmp_path / "dashboard-latest-export-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Dashboard Latest Export",
        "description": "Verify dashboard latest_export stays pinned to the newest run.",
        "source_mode": "file",
        "tags": ["phase7", "dashboard"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "dashboard-latest-export-checksum",
        "extracted_summary": "A compact source summary for dashboard latest export tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief_id = brief_response.json()["brief"]["id"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief_id})
        assert outline_response.status_code == 200
        outline_id = outline_response.json()["outline"]["id"]

        slide_plan_response = client.post(f"/projects/{project_id}/slide-plan:generate", json={"outline_id": outline_id})
        assert slide_plan_response.status_code == 200
        slide_plan_id = slide_plan_response.json()["slide_plan"]["id"]

        artifact_response = client.post(f"/projects/{project_id}/artifact:generate", json={"slide_plan_id": slide_plan_id})
        assert artifact_response.status_code == 200
        artifact_id = artifact_response.json()["artifact"]["id"]

        first_export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact_id, "export_format": "pptx"},
        )
        assert first_export_response.status_code == 200
        first_export_job = first_export_response.json()["export_job"]

        second_export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact_id, "export_format": "pdf"},
        )
        assert second_export_response.status_code == 200
        second_export_job = second_export_response.json()["export_job"]

        history_response = client.get(f"/projects/{project_id}/exports?limit=5")
        assert history_response.status_code == 200
        history_payload = history_response.json()

        dashboard_response = client.get("/projects")
        assert dashboard_response.status_code == 200
        dashboard_payload = dashboard_response.json()

    assert len(history_payload["exports"]) == 2
    assert history_payload["exports"][0]["id"] == second_export_job["id"]
    assert history_payload["exports"][1]["id"] == first_export_job["id"]
    assert len(dashboard_payload["projects"]) == 1
    dashboard_item = dashboard_payload["projects"][0]
    assert dashboard_item["project"]["id"] == project_id
    assert dashboard_item["latest_export"]["id"] == second_export_job["id"]
    assert dashboard_item["latest_export"]["run_id"] == second_export_job["run_id"]
    assert dashboard_item["latest_export"]["preview_pdf_path"] == second_export_job["preview_pdf_path"]
    assert dashboard_item["current_task"]["task_type"] == "export"
    assert dashboard_item["current_task"]["task_status"] == "succeeded"


def test_project_status_keeps_recent_tasks_tail_aligned_with_latest_export_run(tmp_path: Path) -> None:
    database_path = tmp_path / "project-status-latest-export-tail.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Status Latest Export Tail",
        "description": "Verify project status keeps current_task aligned with the newest export run.",
        "source_mode": "file",
        "tags": ["phase7", "status"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "status-latest-export-tail-checksum",
        "extracted_summary": "A compact source summary for project status latest export tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief_id = brief_response.json()["brief"]["id"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief_id})
        assert outline_response.status_code == 200
        outline_id = outline_response.json()["outline"]["id"]

        slide_plan_response = client.post(f"/projects/{project_id}/slide-plan:generate", json={"outline_id": outline_id})
        assert slide_plan_response.status_code == 200
        slide_plan_id = slide_plan_response.json()["slide_plan"]["id"]

        artifact_response = client.post(f"/projects/{project_id}/artifact:generate", json={"slide_plan_id": slide_plan_id})
        assert artifact_response.status_code == 200
        artifact_id = artifact_response.json()["artifact"]["id"]

        first_export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact_id, "export_format": "pptx"},
        )
        assert first_export_response.status_code == 200
        first_export_job = first_export_response.json()["export_job"]

        second_export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact_id, "export_format": "pdf"},
        )
        assert second_export_response.status_code == 200
        second_export_job = second_export_response.json()["export_job"]

        status_response = client.get(f"/projects/{project_id}/status")
        assert status_response.status_code == 200
        status_payload = status_response.json()

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        export_history_response = client.get(f"/projects/{project_id}/exports?limit=5")
        assert export_history_response.status_code == 200
        export_history_payload = export_history_response.json()

    assert status_payload["project_id"] == project_id
    assert status_payload["project_status"] == "exported"
    assert status_payload["current_task"]["id"] == status_payload["recent_tasks"][-1]["id"]
    assert status_payload["current_task"]["task_type"] == "export"
    assert status_payload["current_task"]["task_status"] == "succeeded"
    assert status_payload["current_task"]["result"]["run_id"] == second_export_job["run_id"]
    assert status_payload["recent_tasks"][-2]["task_type"] == "export"
    assert status_payload["recent_tasks"][-2]["result"]["run_id"] == first_export_job["run_id"]
    assert status_payload["recent_tasks"][-1]["task_type"] == "export"
    assert status_payload["recent_tasks"][-1]["result"]["run_id"] == second_export_job["run_id"]
    assert detail_payload["latest_export"]["id"] == second_export_job["id"]
    assert detail_payload["latest_export"]["run_id"] == status_payload["current_task"]["result"]["run_id"]
    assert export_history_payload["exports"][0]["id"] == second_export_job["id"]
    assert export_history_payload["exports"][0]["run_id"] == status_payload["recent_tasks"][-1]["result"]["run_id"]


def test_list_projects_dashboard_keeps_multi_project_order_and_read_models_isolated(tmp_path: Path) -> None:
    database_path = tmp_path / "dashboard-multi-project-order.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    first_project_payload = {
        "name": "Dashboard First Project",
        "description": "Verify dashboard preserves per-project aggregation boundaries.",
        "source_mode": "file",
        "tags": ["phase7", "dashboard"],
    }
    second_project_payload = {
        "name": "Dashboard Second Project",
        "description": "Verify dashboard ordering follows newest project first.",
        "source_mode": "file",
        "tags": ["phase7", "dashboard"],
    }
    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "dashboard-multi-project-checksum",
        "extracted_summary": "A compact source summary for dashboard multi-project tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        first_create_response = client.post("/projects", json=first_project_payload)
        assert first_create_response.status_code == 201
        first_project_id = first_create_response.json()["project"]["id"]

        first_register_response = client.post(f"/projects/{first_project_id}/files", json=file_payload)
        assert first_register_response.status_code == 201

        first_brief_response = client.post(f"/projects/{first_project_id}/brief:generate", json={"force_regenerate": True})
        assert first_brief_response.status_code == 200
        first_brief_id = first_brief_response.json()["brief"]["id"]

        first_outline_response = client.post(f"/projects/{first_project_id}/outline:generate", json={"brief_id": first_brief_id})
        assert first_outline_response.status_code == 200
        first_outline_id = first_outline_response.json()["outline"]["id"]

        first_slide_plan_response = client.post(f"/projects/{first_project_id}/slide-plan:generate", json={"outline_id": first_outline_id})
        assert first_slide_plan_response.status_code == 200
        first_slide_plan_id = first_slide_plan_response.json()["slide_plan"]["id"]

        first_artifact_response = client.post(f"/projects/{first_project_id}/artifact:generate", json={"slide_plan_id": first_slide_plan_id})
        assert first_artifact_response.status_code == 200
        first_artifact_id = first_artifact_response.json()["artifact"]["id"]

        first_export_response = client.post(
            f"/projects/{first_project_id}/export",
            json={"artifact_id": first_artifact_id, "export_format": "pptx"},
        )
        assert first_export_response.status_code == 200
        first_export_job = first_export_response.json()["export_job"]

        second_create_response = client.post("/projects", json=second_project_payload)
        assert second_create_response.status_code == 201
        second_project_id = second_create_response.json()["project"]["id"]

        second_register_response = client.post(f"/projects/{second_project_id}/files", json=file_payload)
        assert second_register_response.status_code == 201

        second_brief_response = client.post(f"/projects/{second_project_id}/brief:generate", json={"force_regenerate": True})
        assert second_brief_response.status_code == 200
        second_brief_id = second_brief_response.json()["brief"]["id"]

        second_outline_response = client.post(f"/projects/{second_project_id}/outline:generate", json={"brief_id": second_brief_id})
        assert second_outline_response.status_code == 200
        second_outline_id = second_outline_response.json()["outline"]["id"]

        second_slide_plan_response = client.post(f"/projects/{second_project_id}/slide-plan:generate", json={"outline_id": second_outline_id})
        assert second_slide_plan_response.status_code == 200
        second_slide_plan_id = second_slide_plan_response.json()["slide_plan"]["id"]

        second_artifact_response = client.post(f"/projects/{second_project_id}/artifact:generate", json={"slide_plan_id": second_slide_plan_id})
        assert second_artifact_response.status_code == 200
        second_artifact_id = second_artifact_response.json()["artifact"]["id"]

        second_export_response = client.post(
            f"/projects/{second_project_id}/export",
            json={"artifact_id": second_artifact_id, "export_format": "pdf"},
        )
        assert second_export_response.status_code == 200
        second_export_job = second_export_response.json()["export_job"]

        dashboard_response = client.get("/projects")
        assert dashboard_response.status_code == 200
        dashboard_payload = dashboard_response.json()

    assert len(dashboard_payload["projects"]) == 2
    first_dashboard_item = dashboard_payload["projects"][0]
    second_dashboard_item = dashboard_payload["projects"][1]

    assert first_dashboard_item["project"]["id"] == second_project_id
    assert first_dashboard_item["project"]["name"] == second_project_payload["name"]
    assert first_dashboard_item["latest_export"]["id"] == second_export_job["id"]
    assert first_dashboard_item["latest_export"]["run_id"] == second_export_job["run_id"]
    assert first_dashboard_item["current_task"]["project_id"] == second_project_id
    assert first_dashboard_item["current_task"]["task_type"] == "export"

    assert second_dashboard_item["project"]["id"] == first_project_id
    assert second_dashboard_item["project"]["name"] == first_project_payload["name"]
    assert second_dashboard_item["latest_export"]["id"] == first_export_job["id"]
    assert second_dashboard_item["latest_export"]["run_id"] == first_export_job["run_id"]
    assert second_dashboard_item["current_task"]["project_id"] == first_project_id
    assert second_dashboard_item["current_task"]["task_type"] == "export"


def test_list_projects_dashboard_keeps_mixed_project_statuses_aligned(tmp_path: Path) -> None:
    database_path = tmp_path / "dashboard-mixed-project-statuses.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    rendered_project_payload = {
        "name": "Dashboard Rendered Project",
        "description": "Verify dashboard keeps rendered project state without export data.",
        "source_mode": "file",
        "tags": ["phase7", "dashboard"],
    }
    exported_project_payload = {
        "name": "Dashboard Exported Project",
        "description": "Verify dashboard keeps exported project state with latest export data.",
        "source_mode": "file",
        "tags": ["phase7", "dashboard"],
    }
    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "dashboard-mixed-status-checksum",
        "extracted_summary": "A compact source summary for dashboard mixed status tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        rendered_create_response = client.post("/projects", json=rendered_project_payload)
        assert rendered_create_response.status_code == 201
        rendered_project_id = rendered_create_response.json()["project"]["id"]

        rendered_register_response = client.post(f"/projects/{rendered_project_id}/files", json=file_payload)
        assert rendered_register_response.status_code == 201

        rendered_brief_response = client.post(f"/projects/{rendered_project_id}/brief:generate", json={"force_regenerate": True})
        assert rendered_brief_response.status_code == 200
        rendered_brief_id = rendered_brief_response.json()["brief"]["id"]

        rendered_outline_response = client.post(f"/projects/{rendered_project_id}/outline:generate", json={"brief_id": rendered_brief_id})
        assert rendered_outline_response.status_code == 200
        rendered_outline_id = rendered_outline_response.json()["outline"]["id"]

        rendered_slide_plan_response = client.post(f"/projects/{rendered_project_id}/slide-plan:generate", json={"outline_id": rendered_outline_id})
        assert rendered_slide_plan_response.status_code == 200
        rendered_slide_plan_id = rendered_slide_plan_response.json()["slide_plan"]["id"]

        rendered_artifact_response = client.post(f"/projects/{rendered_project_id}/artifact:generate", json={"slide_plan_id": rendered_slide_plan_id})
        assert rendered_artifact_response.status_code == 200
        rendered_artifact = rendered_artifact_response.json()["artifact"]

        exported_create_response = client.post("/projects", json=exported_project_payload)
        assert exported_create_response.status_code == 201
        exported_project_id = exported_create_response.json()["project"]["id"]

        exported_register_response = client.post(f"/projects/{exported_project_id}/files", json=file_payload)
        assert exported_register_response.status_code == 201

        exported_brief_response = client.post(f"/projects/{exported_project_id}/brief:generate", json={"force_regenerate": True})
        assert exported_brief_response.status_code == 200
        exported_brief_id = exported_brief_response.json()["brief"]["id"]

        exported_outline_response = client.post(f"/projects/{exported_project_id}/outline:generate", json={"brief_id": exported_brief_id})
        assert exported_outline_response.status_code == 200
        exported_outline_id = exported_outline_response.json()["outline"]["id"]

        exported_slide_plan_response = client.post(f"/projects/{exported_project_id}/slide-plan:generate", json={"outline_id": exported_outline_id})
        assert exported_slide_plan_response.status_code == 200
        exported_slide_plan_id = exported_slide_plan_response.json()["slide_plan"]["id"]

        exported_artifact_response = client.post(f"/projects/{exported_project_id}/artifact:generate", json={"slide_plan_id": exported_slide_plan_id})
        assert exported_artifact_response.status_code == 200
        exported_artifact_id = exported_artifact_response.json()["artifact"]["id"]

        export_response = client.post(
            f"/projects/{exported_project_id}/export",
            json={"artifact_id": exported_artifact_id, "export_format": "pdf"},
        )
        assert export_response.status_code == 200
        export_job = export_response.json()["export_job"]

        dashboard_response = client.get("/projects")
        assert dashboard_response.status_code == 200
        dashboard_payload = dashboard_response.json()

    assert len(dashboard_payload["projects"]) == 2
    newest_dashboard_item = dashboard_payload["projects"][0]
    older_dashboard_item = dashboard_payload["projects"][1]

    assert newest_dashboard_item["project"]["id"] == exported_project_id
    assert newest_dashboard_item["project"]["status"] == "exported"
    assert newest_dashboard_item["latest_artifact"]["id"] == exported_artifact_id
    assert newest_dashboard_item["latest_export"]["id"] == export_job["id"]
    assert newest_dashboard_item["current_task"]["task_type"] == "export"
    assert newest_dashboard_item["current_task"]["project_id"] == exported_project_id

    assert older_dashboard_item["project"]["id"] == rendered_project_id
    assert older_dashboard_item["project"]["status"] == "finalized"
    assert older_dashboard_item["latest_artifact"]["id"] == rendered_artifact["id"]
    assert older_dashboard_item["latest_export"] is None
    assert older_dashboard_item["current_task"]["task_type"] == "render"
    assert older_dashboard_item["current_task"]["project_id"] == rendered_project_id


def test_list_projects_dashboard_keeps_created_rendered_and_exported_states_isolated(tmp_path: Path) -> None:
    database_path = tmp_path / "dashboard-created-rendered-exported.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    created_project_payload = {
        "name": "Dashboard Created Project",
        "description": "Verify dashboard keeps untouched projects empty and isolated.",
        "source_mode": "chat",
        "tags": ["phase7", "dashboard"],
    }
    rendered_project_payload = {
        "name": "Dashboard Rendered Middle Project",
        "description": "Verify dashboard keeps rendered projects without export data.",
        "source_mode": "file",
        "tags": ["phase7", "dashboard"],
    }
    exported_project_payload = {
        "name": "Dashboard Exported Newest Project",
        "description": "Verify dashboard keeps exported projects with latest export data.",
        "source_mode": "file",
        "tags": ["phase7", "dashboard"],
    }
    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "dashboard-created-rendered-exported-checksum",
        "extracted_summary": "A compact source summary for dashboard tri-state tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        created_response = client.post("/projects", json=created_project_payload)
        assert created_response.status_code == 201
        created_project_id = created_response.json()["project"]["id"]

        rendered_create_response = client.post("/projects", json=rendered_project_payload)
        assert rendered_create_response.status_code == 201
        rendered_project_id = rendered_create_response.json()["project"]["id"]

        rendered_register_response = client.post(f"/projects/{rendered_project_id}/files", json=file_payload)
        assert rendered_register_response.status_code == 201

        rendered_brief_response = client.post(f"/projects/{rendered_project_id}/brief:generate", json={"force_regenerate": True})
        assert rendered_brief_response.status_code == 200
        rendered_brief = rendered_brief_response.json()["brief"]

        rendered_outline_response = client.post(f"/projects/{rendered_project_id}/outline:generate", json={"brief_id": rendered_brief["id"]})
        assert rendered_outline_response.status_code == 200
        rendered_outline = rendered_outline_response.json()["outline"]

        rendered_slide_plan_response = client.post(
            f"/projects/{rendered_project_id}/slide-plan:generate",
            json={"outline_id": rendered_outline["id"]},
        )
        assert rendered_slide_plan_response.status_code == 200
        rendered_slide_plan = rendered_slide_plan_response.json()["slide_plan"]

        rendered_artifact_response = client.post(
            f"/projects/{rendered_project_id}/artifact:generate",
            json={"slide_plan_id": rendered_slide_plan["id"]},
        )
        assert rendered_artifact_response.status_code == 200
        rendered_artifact = rendered_artifact_response.json()["artifact"]

        exported_create_response = client.post("/projects", json=exported_project_payload)
        assert exported_create_response.status_code == 201
        exported_project_id = exported_create_response.json()["project"]["id"]

        exported_register_response = client.post(f"/projects/{exported_project_id}/files", json=file_payload)
        assert exported_register_response.status_code == 201

        exported_brief_response = client.post(f"/projects/{exported_project_id}/brief:generate", json={"force_regenerate": True})
        assert exported_brief_response.status_code == 200
        exported_brief = exported_brief_response.json()["brief"]

        exported_outline_response = client.post(f"/projects/{exported_project_id}/outline:generate", json={"brief_id": exported_brief["id"]})
        assert exported_outline_response.status_code == 200
        exported_outline = exported_outline_response.json()["outline"]

        exported_slide_plan_response = client.post(
            f"/projects/{exported_project_id}/slide-plan:generate",
            json={"outline_id": exported_outline["id"]},
        )
        assert exported_slide_plan_response.status_code == 200
        exported_slide_plan = exported_slide_plan_response.json()["slide_plan"]

        exported_artifact_response = client.post(
            f"/projects/{exported_project_id}/artifact:generate",
            json={"slide_plan_id": exported_slide_plan["id"]},
        )
        assert exported_artifact_response.status_code == 200
        exported_artifact = exported_artifact_response.json()["artifact"]

        export_response = client.post(
            f"/projects/{exported_project_id}/export",
            json={"artifact_id": exported_artifact["id"], "export_format": "pptx"},
        )
        assert export_response.status_code == 200
        export_job = export_response.json()["export_job"]

        dashboard_response = client.get("/projects")
        assert dashboard_response.status_code == 200
        dashboard_payload = dashboard_response.json()

    assert len(dashboard_payload["projects"]) == 3

    newest_dashboard_item = dashboard_payload["projects"][0]
    middle_dashboard_item = dashboard_payload["projects"][1]
    oldest_dashboard_item = dashboard_payload["projects"][2]

    assert newest_dashboard_item["project"]["id"] == exported_project_id
    assert newest_dashboard_item["project"]["status"] == "exported"
    assert newest_dashboard_item["latest_brief"]["id"] == exported_brief["id"]
    assert newest_dashboard_item["latest_outline"]["id"] == exported_outline["id"]
    assert newest_dashboard_item["latest_slide_plan"]["id"] == exported_slide_plan["id"]
    assert newest_dashboard_item["latest_artifact"]["id"] == exported_artifact["id"]
    assert newest_dashboard_item["latest_export"]["id"] == export_job["id"]
    assert newest_dashboard_item["current_task"]["task_type"] == "export"
    assert newest_dashboard_item["current_task"]["project_id"] == exported_project_id

    assert middle_dashboard_item["project"]["id"] == rendered_project_id
    assert middle_dashboard_item["project"]["status"] == "finalized"
    assert middle_dashboard_item["latest_brief"]["id"] == rendered_brief["id"]
    assert middle_dashboard_item["latest_outline"]["id"] == rendered_outline["id"]
    assert middle_dashboard_item["latest_slide_plan"]["id"] == rendered_slide_plan["id"]
    assert middle_dashboard_item["latest_artifact"]["id"] == rendered_artifact["id"]
    assert middle_dashboard_item["latest_export"] is None
    assert middle_dashboard_item["current_task"]["task_type"] == "render"
    assert middle_dashboard_item["current_task"]["project_id"] == rendered_project_id

    assert oldest_dashboard_item["project"]["id"] == created_project_id
    assert oldest_dashboard_item["project"]["status"] == "created"
    assert oldest_dashboard_item["latest_brief"] is None
    assert oldest_dashboard_item["latest_outline"] is None
    assert oldest_dashboard_item["latest_slide_plan"] is None
    assert oldest_dashboard_item["latest_artifact"] is None
    assert oldest_dashboard_item["latest_export"] is None
    assert oldest_dashboard_item["current_task"]["project_id"] == created_project_id
    assert oldest_dashboard_item["current_task"]["task_status"] in {"pending", "succeeded"}


def test_list_projects_dashboard_keeps_export_failed_projects_isolated(tmp_path: Path) -> None:
    database_path = tmp_path / "dashboard-export-failed-projects.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    failed_project_payload = {
        "name": "Dashboard Failed Export Project",
        "description": "Verify dashboard keeps export_failed projects from inheriting export data.",
        "source_mode": "file",
        "tags": ["phase7", "dashboard"],
    }
    exported_project_payload = {
        "name": "Dashboard Healthy Export Project",
        "description": "Verify dashboard keeps exported projects separate from failed export projects.",
        "source_mode": "file",
        "tags": ["phase7", "dashboard"],
    }
    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "dashboard-export-failed-checksum",
        "extracted_summary": "A compact source summary for dashboard export failure tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        failed_create_response = client.post("/projects", json=failed_project_payload)
        assert failed_create_response.status_code == 201
        failed_project_id = failed_create_response.json()["project"]["id"]

        failed_export_response = client.post(
            f"/projects/{failed_project_id}/export",
            json={"artifact_id": "missing-artifact", "export_format": "pptx"},
        )
        assert failed_export_response.status_code == 400

        exported_create_response = client.post("/projects", json=exported_project_payload)
        assert exported_create_response.status_code == 201
        exported_project_id = exported_create_response.json()["project"]["id"]

        exported_register_response = client.post(f"/projects/{exported_project_id}/files", json=file_payload)
        assert exported_register_response.status_code == 201

        exported_brief_response = client.post(f"/projects/{exported_project_id}/brief:generate", json={"force_regenerate": True})
        assert exported_brief_response.status_code == 200
        exported_brief = exported_brief_response.json()["brief"]

        exported_outline_response = client.post(f"/projects/{exported_project_id}/outline:generate", json={"brief_id": exported_brief["id"]})
        assert exported_outline_response.status_code == 200
        exported_outline = exported_outline_response.json()["outline"]

        exported_slide_plan_response = client.post(
            f"/projects/{exported_project_id}/slide-plan:generate",
            json={"outline_id": exported_outline["id"]},
        )
        assert exported_slide_plan_response.status_code == 200
        exported_slide_plan = exported_slide_plan_response.json()["slide_plan"]

        exported_artifact_response = client.post(
            f"/projects/{exported_project_id}/artifact:generate",
            json={"slide_plan_id": exported_slide_plan["id"]},
        )
        assert exported_artifact_response.status_code == 200
        exported_artifact = exported_artifact_response.json()["artifact"]

        exported_export_response = client.post(
            f"/projects/{exported_project_id}/export",
            json={"artifact_id": exported_artifact["id"], "export_format": "pdf"},
        )
        assert exported_export_response.status_code == 200
        exported_export_job = exported_export_response.json()["export_job"]

        dashboard_response = client.get("/projects")
        assert dashboard_response.status_code == 200
        dashboard_payload = dashboard_response.json()

    assert len(dashboard_payload["projects"]) == 2

    newest_dashboard_item = dashboard_payload["projects"][0]
    older_dashboard_item = dashboard_payload["projects"][1]

    assert newest_dashboard_item["project"]["id"] == exported_project_id
    assert newest_dashboard_item["project"]["status"] == "exported"
    assert newest_dashboard_item["latest_brief"]["id"] == exported_brief["id"]
    assert newest_dashboard_item["latest_outline"]["id"] == exported_outline["id"]
    assert newest_dashboard_item["latest_slide_plan"]["id"] == exported_slide_plan["id"]
    assert newest_dashboard_item["latest_artifact"]["id"] == exported_artifact["id"]
    assert newest_dashboard_item["latest_export"]["id"] == exported_export_job["id"]
    assert newest_dashboard_item["current_task"]["task_type"] == "export"
    assert newest_dashboard_item["current_task"]["task_status"] == "succeeded"
    assert newest_dashboard_item["current_task"]["project_id"] == exported_project_id

    assert older_dashboard_item["project"]["id"] == failed_project_id
    assert older_dashboard_item["project"]["status"] == "export_failed"
    assert older_dashboard_item["latest_brief"] is None
    assert older_dashboard_item["latest_outline"] is None
    assert older_dashboard_item["latest_slide_plan"] is None
    assert older_dashboard_item["latest_artifact"] is None
    assert older_dashboard_item["latest_export"] is None
    assert older_dashboard_item["current_task"]["task_type"] == "export"
    assert older_dashboard_item["current_task"]["task_status"] == "failed"
    assert older_dashboard_item["current_task"]["project_id"] == failed_project_id
    assert older_dashboard_item["current_task"]["error_message"] == "Project has no rendered artifact to export"


def test_project_status_keeps_failed_export_tail_aligned_with_project_status(tmp_path: Path) -> None:
    database_path = tmp_path / "status-export-failed-tail.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Status Failed Export Project",
        "description": "Verify failed export tail stays aligned across detail and status read models.",
        "source_mode": "file",
        "tags": ["phase7", "status"],
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": "missing-artifact", "export_format": "pptx"},
        )
        assert export_response.status_code == 400
        assert export_response.json()["detail"] == "Project has no rendered artifact to export"

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        status_response = client.get(f"/projects/{project_id}/status")
        assert status_response.status_code == 200
        status_payload = status_response.json()

    assert detail_payload["project"]["status"] == "export_failed"
    assert detail_payload["latest_export"] is None
    assert status_payload["project_status"] == "export_failed"
    assert status_payload["current_task"]["id"] == status_payload["recent_tasks"][-1]["id"]
    assert status_payload["current_task"]["task_type"] == "export"
    assert status_payload["current_task"]["task_status"] == "failed"
    assert status_payload["current_task"]["project_id"] == project_id
    assert status_payload["current_task"]["error_message"] == "Project has no rendered artifact to export"
    assert status_payload["recent_tasks"][-1]["task_type"] == "export"
    assert status_payload["recent_tasks"][-1]["task_status"] == "failed"
    assert status_payload["recent_tasks"][-1]["project_id"] == project_id
    assert status_payload["recent_tasks"][-1]["error_message"] == "Project has no rendered artifact to export"


def test_project_status_keeps_render_tail_aligned_before_any_export(tmp_path: Path) -> None:
    database_path = tmp_path / "status-render-tail.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Status Render Tail Project",
        "description": "Verify rendered projects keep status tail aligned before any export exists.",
        "source_mode": "file",
        "tags": ["phase7", "status"],
    }
    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "status-render-tail-checksum",
        "extracted_summary": "A compact source summary for rendered status tail tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief = brief_response.json()["brief"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief["id"]})
        assert outline_response.status_code == 200
        outline = outline_response.json()["outline"]

        slide_plan_response = client.post(
            f"/projects/{project_id}/slide-plan:generate",
            json={"outline_id": outline["id"]},
        )
        assert slide_plan_response.status_code == 200
        slide_plan = slide_plan_response.json()["slide_plan"]

        artifact_response = client.post(
            f"/projects/{project_id}/artifact:generate",
            json={"slide_plan_id": slide_plan["id"]},
        )
        assert artifact_response.status_code == 200
        artifact = artifact_response.json()["artifact"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        status_response = client.get(f"/projects/{project_id}/status")
        assert status_response.status_code == 200
        status_payload = status_response.json()

        export_history_response = client.get(f"/projects/{project_id}/exports")
        assert export_history_response.status_code == 200
        export_history_payload = export_history_response.json()

    assert detail_payload["project"]["status"] == "finalized"
    assert detail_payload["latest_artifact"]["id"] == artifact["id"]
    assert detail_payload["latest_export"] is None
    assert status_payload["project_status"] == "finalized"
    assert status_payload["current_task"]["id"] == status_payload["recent_tasks"][-1]["id"]
    assert status_payload["current_task"]["task_type"] == "render"
    assert status_payload["current_task"]["task_status"] == "succeeded"
    assert status_payload["current_task"]["project_id"] == project_id
    assert status_payload["current_task"]["result"]["artifact_id"] == artifact["id"]
    assert status_payload["recent_tasks"][-1]["task_type"] == "render"
    assert status_payload["recent_tasks"][-1]["task_status"] == "succeeded"
    assert status_payload["recent_tasks"][-1]["project_id"] == project_id
    assert status_payload["recent_tasks"][-1]["result"]["artifact_id"] == artifact["id"]
    assert export_history_payload["exports"] == []


def test_project_status_keeps_brief_tail_aligned_before_outline(tmp_path: Path) -> None:
    database_path = tmp_path / "status-brief-tail.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Status Brief Tail Project",
        "description": "Verify briefing state stays aligned before outline generation begins.",
        "source_mode": "file",
        "tags": ["phase7", "status"],
    }
    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "status-brief-tail-checksum",
        "extracted_summary": "A compact source summary for briefing state tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief = brief_response.json()["brief"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        status_response = client.get(f"/projects/{project_id}/status")
        assert status_response.status_code == 200
        status_payload = status_response.json()

        export_history_response = client.get(f"/projects/{project_id}/exports")
        assert export_history_response.status_code == 200
        export_history_payload = export_history_response.json()

    assert detail_payload["project"]["status"] == "briefing"
    assert detail_payload["latest_brief"]["id"] == brief["id"]
    assert detail_payload["latest_outline"] is None
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None
    assert status_payload["project_status"] == "briefing"
    assert status_payload["current_task"]["id"] == status_payload["recent_tasks"][-1]["id"]
    assert status_payload["current_task"]["task_type"] == "generate_brief"
    assert status_payload["current_task"]["task_status"] == "succeeded"
    assert status_payload["current_task"]["project_id"] == project_id
    assert status_payload["current_task"]["result"]["brief_id"] == brief["id"]
    assert status_payload["recent_tasks"][-1]["task_type"] == "generate_brief"
    assert status_payload["recent_tasks"][-1]["task_status"] == "succeeded"
    assert status_payload["recent_tasks"][-1]["project_id"] == project_id
    assert status_payload["recent_tasks"][-1]["result"]["brief_id"] == brief["id"]
    assert export_history_payload["exports"] == []


def test_project_status_keeps_analyzed_tail_aligned_before_brief(tmp_path: Path) -> None:
    database_path = tmp_path / "status-analyzed-tail.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Analyzed State Demo",
        "description": "Verify analyzed status remains aligned before any brief exists.",
        "source_mode": "file",
        "tags": ["phase7", "status", "analyzed"],
    }

    upload_payload = {
        "file_name": "analyzed.md",
        "file_type": "markdown",
        "mime_type": "text/markdown",
        "content_base64": base64.b64encode(b"   \n\n   ").decode("ascii"),
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        upload_response = client.post(f"/projects/{project_id}/files:upload", json=upload_payload)
        assert upload_response.status_code == 201
        file_id = upload_response.json()["file"]["id"]

        parse_response = client.post(
            f"/projects/{project_id}/files:parse",
            json={"file_ids": [file_id], "rebuild_bundle": True},
        )
        assert parse_response.status_code == 200
        source_bundle = parse_response.json()["source_bundle"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        status_response = client.get(f"/projects/{project_id}/status")
        assert status_response.status_code == 200
        status_payload = status_response.json()

        exports_response = client.get(f"/projects/{project_id}/exports")
        assert exports_response.status_code == 200
        exports_payload = exports_response.json()

    assert detail_payload["project"]["status"] == "analyzed"
    assert detail_payload["latest_source_bundle"]["id"] == source_bundle["id"]
    assert detail_payload["latest_source_bundle"]["status"] == "needs_review"
    assert detail_payload["latest_source_bundle"]["metadata"]["warnings_by_file"]
    assert detail_payload["latest_brief"] is None
    assert detail_payload["latest_outline"] is None
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None

    assert status_payload["project_status"] == "analyzed"
    assert status_payload["current_task"]["id"] == status_payload["recent_tasks"][-1]["id"]
    assert status_payload["current_task"]["task_type"] == "parse"
    assert status_payload["current_task"]["task_status"] == "succeeded"
    assert status_payload["current_task"]["project_id"] == project_id
    assert status_payload["current_task"]["result"]["parsed_file_count"] == 1
    assert status_payload["current_task"]["result"]["failed_file_count"] == 0
    assert status_payload["current_task"]["result"]["source_bundle_id"] == source_bundle["id"]
    assert status_payload["recent_tasks"][-1]["task_type"] == "parse"
    assert status_payload["recent_tasks"][-1]["task_status"] == "succeeded"
    assert status_payload["recent_tasks"][-1]["project_id"] == project_id
    assert status_payload["recent_tasks"][-1]["result"]["parsed_file_count"] == 1
    assert status_payload["recent_tasks"][-1]["result"]["failed_file_count"] == 0
    assert status_payload["recent_tasks"][-1]["result"]["source_bundle_id"] == source_bundle["id"]
    assert exports_payload["exports"] == []


def test_project_status_keeps_parse_failed_tail_aligned_before_brief(tmp_path: Path) -> None:
    database_path = tmp_path / "status-parse-failed-tail.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Parse Failed State Demo",
        "description": "Verify parse_failed status remains aligned before any brief exists.",
        "source_mode": "file",
        "tags": ["phase7", "status", "parse_failed"],
    }

    upload_payload = {
        "file_name": "broken.url",
        "file_type": "url",
        "mime_type": "text/plain",
        "content_base64": base64.b64encode(b"http://127.0.0.1:1/parse-failure-fixture").decode("ascii"),
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        upload_response = client.post(f"/projects/{project_id}/files:upload", json=upload_payload)
        assert upload_response.status_code == 201
        file_id = upload_response.json()["file"]["id"]

        parse_response = client.post(
            f"/projects/{project_id}/files:parse",
            json={"file_ids": [file_id], "rebuild_bundle": True},
        )
        assert parse_response.status_code == 200
        parse_payload = parse_response.json()

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        status_response = client.get(f"/projects/{project_id}/status")
        assert status_response.status_code == 200
        status_payload = status_response.json()

        exports_response = client.get(f"/projects/{project_id}/exports")
        assert exports_response.status_code == 200
        exports_payload = exports_response.json()

    parsed_file = parse_payload["files"][0]

    assert parse_payload["source_bundle"] is None
    assert parsed_file["id"] == file_id
    assert parsed_file["parse_status"] == "failed"
    assert parsed_file["parse_error"]

    assert detail_payload["project"]["status"] == "parse_failed"
    assert detail_payload["latest_source_bundle"] is None
    assert detail_payload["latest_brief"] is None
    assert detail_payload["latest_outline"] is None
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None

    assert status_payload["project_status"] == "parse_failed"
    assert status_payload["current_task"]["id"] == status_payload["recent_tasks"][-1]["id"]
    assert status_payload["current_task"]["task_type"] == "parse"
    assert status_payload["current_task"]["task_status"] == "failed"
    assert status_payload["current_task"]["project_id"] == project_id
    assert status_payload["current_task"]["error_message"] == "Some files failed during parsing"
    assert status_payload["current_task"]["result"]["parsed_file_count"] == 0
    assert status_payload["current_task"]["result"]["failed_file_count"] == 1
    assert status_payload["current_task"]["result"]["source_bundle_id"] is None
    assert status_payload["recent_tasks"][-1]["task_type"] == "parse"
    assert status_payload["recent_tasks"][-1]["task_status"] == "failed"
    assert status_payload["recent_tasks"][-1]["project_id"] == project_id
    assert status_payload["recent_tasks"][-1]["error_message"] == "Some files failed during parsing"
    assert status_payload["recent_tasks"][-1]["result"]["parsed_file_count"] == 0
    assert status_payload["recent_tasks"][-1]["result"]["failed_file_count"] == 1
    assert status_payload["recent_tasks"][-1]["result"]["source_bundle_id"] is None
    assert exports_payload["exports"] == []


def test_project_status_keeps_parse_tail_aligned_before_brief(tmp_path: Path) -> None:
    database_path = tmp_path / "status-parsed-tail.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Parsed State Demo",
        "description": "Verify parsed status remains aligned before any brief exists.",
        "source_mode": "file",
        "tags": ["phase7", "status", "parsed"],
    }

    upload_payload = {
        "file_name": "parsed.md",
        "file_type": "markdown",
        "mime_type": "text/markdown",
        "content_base64": base64.b64encode(b"# Parsed Demo\n\nThis file exists on disk so parse can succeed.").decode("ascii"),
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files:upload", json=upload_payload)
        assert register_response.status_code == 201

        parse_response = client.post(
            f"/projects/{project_id}/files:parse",
            json={"file_ids": [register_response.json()["file"]["id"]], "rebuild_bundle": True},
        )
        assert parse_response.status_code == 200
        source_bundle = parse_response.json()["source_bundle"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        status_response = client.get(f"/projects/{project_id}/status")
        assert status_response.status_code == 200
        status_payload = status_response.json()

        exports_response = client.get(f"/projects/{project_id}/exports")
        assert exports_response.status_code == 200
        exports_payload = exports_response.json()

    assert detail_payload["project"]["status"] == "parsed"
    assert detail_payload["latest_source_bundle"]["id"] == source_bundle["id"]
    assert detail_payload["latest_brief"] is None
    assert detail_payload["latest_outline"] is None
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None

    assert status_payload["project_status"] == "parsed"
    assert status_payload["current_task"]["id"] == status_payload["recent_tasks"][-1]["id"]
    assert status_payload["current_task"]["task_type"] == "parse"
    assert status_payload["current_task"]["task_status"] == "succeeded"
    assert status_payload["current_task"]["project_id"] == project_id
    assert status_payload["current_task"]["result"]["parsed_file_count"] == 1
    assert status_payload["current_task"]["result"]["failed_file_count"] == 0
    assert status_payload["current_task"]["result"]["source_bundle_id"] == source_bundle["id"]
    assert status_payload["recent_tasks"][-1]["task_type"] == "parse"
    assert status_payload["recent_tasks"][-1]["task_status"] == "succeeded"
    assert status_payload["recent_tasks"][-1]["project_id"] == project_id
    assert status_payload["recent_tasks"][-1]["result"]["parsed_file_count"] == 1
    assert status_payload["recent_tasks"][-1]["result"]["failed_file_count"] == 0
    assert status_payload["recent_tasks"][-1]["result"]["source_bundle_id"] == source_bundle["id"]
    assert exports_payload["exports"] == []


def test_project_detail_status_and_export_history_remain_consistent_across_pipeline(tmp_path: Path) -> None:
    database_path = tmp_path / "project-detail-status-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Project Detail Status Demo",
        "description": "Verify project detail, status, and export history stay aligned.",
        "source_mode": "file",
        "tags": ["phase7", "api-integration"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "detail-status-checksum",
        "extracted_summary": "A compact source summary for project detail and status tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        initial_detail_response = client.get(f"/projects/{project_id}")
        assert initial_detail_response.status_code == 200
        initial_detail = initial_detail_response.json()

        initial_status_response = client.get(f"/projects/{project_id}/status")
        assert initial_status_response.status_code == 200
        initial_status = initial_status_response.json()

        empty_history_response = client.get(f"/projects/{project_id}/exports")
        assert empty_history_response.status_code == 200
        empty_history = empty_history_response.json()

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief = brief_response.json()["brief"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief["id"]})
        assert outline_response.status_code == 200
        outline = outline_response.json()["outline"]

        slide_plan_response = client.post(f"/projects/{project_id}/slide-plan:generate", json={"outline_id": outline["id"]})
        assert slide_plan_response.status_code == 200
        slide_plan = slide_plan_response.json()["slide_plan"]

        artifact_response = client.post(f"/projects/{project_id}/artifact:generate", json={"slide_plan_id": slide_plan["id"]})
        assert artifact_response.status_code == 200
        artifact = artifact_response.json()["artifact"]

        detail_after_render_response = client.get(f"/projects/{project_id}")
        assert detail_after_render_response.status_code == 200
        detail_after_render = detail_after_render_response.json()

        status_after_render_response = client.get(f"/projects/{project_id}/status")
        assert status_after_render_response.status_code == 200
        status_after_render = status_after_render_response.json()

        export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact["id"], "export_format": "pptx"},
        )
        assert export_response.status_code == 200
        export_job = export_response.json()["export_job"]

        final_detail_response = client.get(f"/projects/{project_id}")
        assert final_detail_response.status_code == 200
        final_detail = final_detail_response.json()

        final_status_response = client.get(f"/projects/{project_id}/status")
        assert final_status_response.status_code == 200
        final_status = final_status_response.json()

        export_history_response = client.get(f"/projects/{project_id}/exports")
        assert export_history_response.status_code == 200
        export_history = export_history_response.json()

        export_detail_response = client.get(f"/projects/{project_id}/exports/{export_job['id']}")
        assert export_detail_response.status_code == 200
        export_detail = export_detail_response.json()

    assert initial_detail["project"]["id"] == project_id
    assert initial_detail["project"]["status"] == "created"
    assert initial_detail["latest_brief"] is None
    assert initial_detail["latest_outline"] is None
    assert initial_detail["latest_slide_plan"] is None
    assert initial_detail["latest_artifact"] is None
    assert initial_detail["latest_export"] is None

    assert initial_status["project_id"] == project_id
    assert initial_status["project_status"] == "created"
    assert initial_status["current_task"]["task_type"] == "ingest"
    assert initial_status["current_task"]["task_status"] == "succeeded"
    assert initial_status["recent_tasks"][-1]["task_type"] == "ingest"
    assert empty_history["exports"] == []

    assert detail_after_render["project"]["status"] == "finalized"
    assert detail_after_render["latest_brief"]["id"] == brief["id"]
    assert detail_after_render["latest_outline"]["id"] == outline["id"]
    assert detail_after_render["latest_slide_plan"]["id"] == slide_plan["id"]
    assert detail_after_render["latest_artifact"]["id"] == artifact["id"]
    assert detail_after_render["latest_export"] is None

    assert status_after_render["project_status"] == "finalized"
    assert status_after_render["current_task"]["task_type"] == "render"
    assert status_after_render["current_task"]["task_status"] == "succeeded"
    assert status_after_render["recent_tasks"][-1]["task_type"] == "render"

    assert final_detail["project"]["status"] == "exported"
    assert final_detail["latest_brief"]["id"] == brief["id"]
    assert final_detail["latest_outline"]["id"] == outline["id"]
    assert final_detail["latest_slide_plan"]["id"] == slide_plan["id"]
    assert final_detail["latest_artifact"]["id"] == artifact["id"]
    assert final_detail["latest_export"]["id"] == export_job["id"]
    assert final_detail["latest_export"]["run_id"] == export_job["run_id"]

    assert final_status["project_status"] == "exported"
    assert final_status["current_task"]["task_type"] == "export"
    assert final_status["current_task"]["task_status"] == "succeeded"
    assert final_status["recent_tasks"][-1]["task_type"] == "export"
    assert len(final_status["recent_tasks"]) >= 5

    assert len(export_history["exports"]) == 1
    assert export_history["exports"][0]["id"] == export_job["id"]
    assert export_history["exports"][0]["run_id"] == export_job["run_id"]
    assert export_detail["export_job"]["id"] == export_job["id"]
    assert export_detail["export_job"]["project_id"] == project_id


def test_project_detail_status_and_export_read_models_return_404_for_missing_project(tmp_path: Path) -> None:
    database_path = tmp_path / "project-read-model-missing.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    with TestClient(app) as client:
        detail_response = client.get("/projects/missing-project")
        status_response = client.get("/projects/missing-project/status")
        exports_response = client.get("/projects/missing-project/exports")
        export_detail_response = client.get("/projects/missing-project/exports/missing-export")

    for response in (detail_response, status_response, exports_response, export_detail_response):
        assert response.status_code == 404
        assert response.json() == {"detail": "Project not found"}


def test_project_status_keeps_created_tail_aligned_before_any_file_registration(tmp_path: Path) -> None:
    database_path = tmp_path / "project-status-created.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Created State Demo",
        "description": "Verify created status remains aligned before any file exists.",
        "source_mode": "chat",
        "tags": ["phase7", "status", "created"],
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        status_response = client.get(f"/projects/{project_id}/status")
        assert status_response.status_code == 200
        status_payload = status_response.json()

        exports_response = client.get(f"/projects/{project_id}/exports")
        assert exports_response.status_code == 200
        exports_payload = exports_response.json()

    assert detail_payload["project"]["status"] == "created"
    assert detail_payload["latest_source_bundle"] is None
    assert detail_payload["latest_brief"] is None
    assert detail_payload["latest_outline"] is None
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None

    assert status_payload["project_status"] == "created"
    assert status_payload["current_task"]["id"] == status_payload["recent_tasks"][-1]["id"]
    assert status_payload["current_task"]["task_type"] == "ingest"
    assert status_payload["current_task"]["task_status"] == "succeeded"
    assert status_payload["current_task"]["project_id"] == project_id
    assert status_payload["recent_tasks"][-1]["task_type"] == "ingest"
    assert status_payload["recent_tasks"][-1]["task_status"] == "succeeded"
    assert status_payload["recent_tasks"][-1]["project_id"] == project_id
    assert exports_payload["exports"] == []


def test_project_detail_keeps_outline_tail_aligned_before_slide_plan(tmp_path: Path) -> None:
    database_path = tmp_path / "project-detail-outlined.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Outlined Detail Demo",
        "description": "Verify outlined detail stays aligned before any slide plan exists.",
        "source_mode": "file",
        "tags": ["phase7", "detail", "outlined"],
    }

    file_payload = {
        "file_name": "outlined-detail.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/outlined-detail.md",
        "mime_type": "text/markdown",
        "size_bytes": 208,
        "checksum": "outlined-detail-checksum",
        "extracted_summary": "A source summary that is sufficient to generate a brief and outline but not a slide plan.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief = brief_response.json()["brief"]

        outline_response = client.post(
            f"/projects/{project_id}/outline:generate",
            json={"brief_id": brief["id"]},
        )
        assert outline_response.status_code == 200
        outline = outline_response.json()["outline"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

    assert detail_payload["project"]["id"] == project_id
    assert detail_payload["project"]["status"] == "outlined"
    assert detail_payload["latest_source_bundle"] is not None
    assert detail_payload["latest_brief"]["id"] == brief["id"]
    assert detail_payload["latest_outline"]["id"] == outline["id"]
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None


def test_project_detail_keeps_brief_tail_aligned_before_outline(tmp_path: Path) -> None:
    database_path = tmp_path / "project-detail-briefing.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Briefing Detail Demo",
        "description": "Verify briefing detail stays aligned before any outline exists.",
        "source_mode": "file",
        "tags": ["phase7", "detail", "briefing"],
    }

    file_payload = {
        "file_name": "briefing-detail.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/briefing-detail.md",
        "mime_type": "text/markdown",
        "size_bytes": 176,
        "checksum": "briefing-detail-checksum",
        "extracted_summary": "A source summary that is sufficient to generate a brief but not yet any later planning objects.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief = brief_response.json()["brief"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

    assert detail_payload["project"]["id"] == project_id
    assert detail_payload["project"]["status"] == "briefing"
    assert detail_payload["latest_source_bundle"] is not None
    assert detail_payload["latest_brief"]["id"] == brief["id"]
    assert detail_payload["latest_outline"] is None
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None


def test_project_detail_keeps_parse_tail_aligned_before_brief(tmp_path: Path) -> None:
    database_path = tmp_path / "project-detail-parsed.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Parsed Detail Demo",
        "description": "Verify parsed detail stays aligned before any brief exists.",
        "source_mode": "file",
        "tags": ["phase7", "detail", "parsed"],
    }

    upload_payload = {
        "file_name": "parsed-detail.md",
        "file_type": "markdown",
        "mime_type": "text/markdown",
        "content_base64": base64.b64encode(b"# Parsed Detail\n\nA minimal markdown file for parsed detail regression.").decode("ascii"),
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        upload_response = client.post(f"/projects/{project_id}/files:upload", json=upload_payload)
        assert upload_response.status_code == 201
        file_id = upload_response.json()["file"]["id"]

        parse_response = client.post(
            f"/projects/{project_id}/files:parse",
            json={"file_ids": [file_id], "rebuild_bundle": True},
        )
        assert parse_response.status_code == 200
        source_bundle = parse_response.json()["source_bundle"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

    assert detail_payload["project"]["id"] == project_id
    assert detail_payload["project"]["status"] == "parsed"
    assert detail_payload["latest_source_bundle"]["id"] == source_bundle["id"]
    assert detail_payload["latest_brief"] is None
    assert detail_payload["latest_outline"] is None
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None


def test_project_detail_keeps_analyzed_tail_aligned_before_brief(tmp_path: Path) -> None:
    database_path = tmp_path / "project-detail-analyzed.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Analyzed Detail Demo",
        "description": "Verify analyzed detail stays aligned before any brief exists.",
        "source_mode": "file",
        "tags": ["phase7", "detail", "analyzed"],
    }

    upload_payload = {
        "file_name": "analyzed-detail.md",
        "file_type": "markdown",
        "mime_type": "text/markdown",
        "content_base64": base64.b64encode(b"   \n\n   ").decode("ascii"),
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        upload_response = client.post(f"/projects/{project_id}/files:upload", json=upload_payload)
        assert upload_response.status_code == 201
        file_id = upload_response.json()["file"]["id"]

        parse_response = client.post(
            f"/projects/{project_id}/files:parse",
            json={"file_ids": [file_id], "rebuild_bundle": True},
        )
        assert parse_response.status_code == 200
        source_bundle = parse_response.json()["source_bundle"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

    assert detail_payload["project"]["id"] == project_id
    assert detail_payload["project"]["status"] == "analyzed"
    assert detail_payload["latest_source_bundle"]["id"] == source_bundle["id"]
    assert detail_payload["latest_source_bundle"]["status"] == "needs_review"
    assert detail_payload["latest_source_bundle"]["metadata"]["warnings_by_file"]
    assert detail_payload["latest_brief"] is None
    assert detail_payload["latest_outline"] is None
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None


def test_project_detail_keeps_parse_failed_tail_aligned_before_brief(tmp_path: Path) -> None:
    database_path = tmp_path / "project-detail-parse-failed.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Parse Failed Detail Demo",
        "description": "Verify parse_failed detail stays aligned before any brief exists.",
        "source_mode": "file",
        "tags": ["phase7", "detail", "parse_failed"],
    }

    upload_payload = {
        "file_name": "broken-detail.url",
        "file_type": "url",
        "mime_type": "text/plain",
        "content_base64": base64.b64encode(b"http://127.0.0.1:1/parse-failure-detail-fixture").decode("ascii"),
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        upload_response = client.post(f"/projects/{project_id}/files:upload", json=upload_payload)
        assert upload_response.status_code == 201
        file_id = upload_response.json()["file"]["id"]

        parse_response = client.post(
            f"/projects/{project_id}/files:parse",
            json={"file_ids": [file_id], "rebuild_bundle": True},
        )
        assert parse_response.status_code == 200
        parse_payload = parse_response.json()

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

    assert parse_payload["source_bundle"] is None
    assert parse_payload["files"][0]["id"] == file_id
    assert parse_payload["files"][0]["parse_status"] == "failed"
    assert parse_payload["files"][0]["parse_error"]

    assert detail_payload["project"]["id"] == project_id
    assert detail_payload["project"]["status"] == "parse_failed"
    assert detail_payload["latest_source_bundle"] is None
    assert detail_payload["latest_brief"] is None
    assert detail_payload["latest_outline"] is None
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None


def test_project_detail_keeps_render_tail_aligned_before_any_export(tmp_path: Path) -> None:
    database_path = tmp_path / "project-detail-finalized.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Finalized Detail Demo",
        "description": "Verify finalized detail stays aligned before any export exists.",
        "source_mode": "file",
        "tags": ["phase7", "detail", "finalized"],
    }

    file_payload = {
        "file_name": "finalized-detail.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/finalized-detail.md",
        "mime_type": "text/markdown",
        "size_bytes": 224,
        "checksum": "finalized-detail-checksum",
        "extracted_summary": "A source summary that is sufficient to render an artifact but not export it yet.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief = brief_response.json()["brief"]

        outline_response = client.post(
            f"/projects/{project_id}/outline:generate",
            json={"brief_id": brief["id"]},
        )
        assert outline_response.status_code == 200
        outline = outline_response.json()["outline"]

        slide_plan_response = client.post(
            f"/projects/{project_id}/slide-plan:generate",
            json={"outline_id": outline["id"]},
        )
        assert slide_plan_response.status_code == 200
        slide_plan = slide_plan_response.json()["slide_plan"]

        artifact_response = client.post(
            f"/projects/{project_id}/artifact:generate",
            json={"slide_plan_id": slide_plan["id"]},
        )
        assert artifact_response.status_code == 200
        artifact = artifact_response.json()["artifact"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

    assert detail_payload["project"]["id"] == project_id
    assert detail_payload["project"]["status"] == "finalized"
    assert detail_payload["latest_brief"]["id"] == brief["id"]
    assert detail_payload["latest_outline"]["id"] == outline["id"]
    assert detail_payload["latest_slide_plan"]["id"] == slide_plan["id"]
    assert detail_payload["latest_artifact"]["id"] == artifact["id"]
    assert detail_payload["latest_export"] is None


def test_project_detail_keeps_failed_export_tail_aligned_with_project_status(tmp_path: Path) -> None:
    database_path = tmp_path / "project-detail-export-failed.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Export Failed Detail Demo",
        "description": "Verify failed export detail stays aligned with project status.",
        "source_mode": "file",
        "tags": ["phase7", "detail", "export_failed"],
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": "missing-artifact", "export_format": "pptx"},
        )
        assert export_response.status_code == 400
        assert export_response.json()["detail"] == "Project has no rendered artifact to export"

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

    assert detail_payload["project"]["id"] == project_id
    assert detail_payload["project"]["status"] == "export_failed"
    assert detail_payload["latest_source_bundle"] is None
    assert detail_payload["latest_brief"] is None
    assert detail_payload["latest_outline"] is None
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None


def test_project_status_keeps_outline_tail_aligned_before_slide_plan(tmp_path: Path) -> None:
    database_path = tmp_path / "project-status-outlined.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Outlined State Demo",
        "description": "Verify outlined status remains aligned before any slide plan exists.",
        "source_mode": "file",
        "tags": ["phase7", "status", "outlined"],
    }

    file_payload = {
        "file_name": "outlined.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/outlined.md",
        "mime_type": "text/markdown",
        "size_bytes": 192,
        "checksum": "outlined-checksum",
        "extracted_summary": "A source summary that is sufficient to generate a brief and outline.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief = brief_response.json()["brief"]

        outline_response = client.post(
            f"/projects/{project_id}/outline:generate",
            json={"brief_id": brief["id"]},
        )
        assert outline_response.status_code == 200
        outline = outline_response.json()["outline"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        status_response = client.get(f"/projects/{project_id}/status")
        assert status_response.status_code == 200
        status_payload = status_response.json()

        exports_response = client.get(f"/projects/{project_id}/exports")
        assert exports_response.status_code == 200
        exports_payload = exports_response.json()

    assert detail_payload["project"]["status"] == "outlined"
    assert detail_payload["latest_brief"]["id"] == brief["id"]
    assert detail_payload["latest_outline"]["id"] == outline["id"]
    assert detail_payload["latest_slide_plan"] is None
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None

    assert status_payload["project_status"] == "outlined"
    assert status_payload["current_task"]["id"] == status_payload["recent_tasks"][-1]["id"]
    assert status_payload["current_task"]["task_type"] == "generate_outline"
    assert status_payload["current_task"]["task_status"] == "succeeded"
    assert status_payload["current_task"]["result"]["outline_id"] == outline["id"]
    assert status_payload["current_task"]["result"]["brief_id"] == brief["id"]
    assert status_payload["recent_tasks"][-1]["task_type"] == "generate_outline"
    assert status_payload["recent_tasks"][-1]["result"]["outline_id"] == outline["id"]
    assert exports_payload["exports"] == []


def test_project_status_keeps_slide_plan_tail_aligned_before_render(tmp_path: Path) -> None:
    database_path = tmp_path / "project-status-planned.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Planned State Demo",
        "description": "Verify planned status remains aligned before any artifact exists.",
        "source_mode": "file",
        "tags": ["phase7", "status", "planned"],
    }

    file_payload = {
        "file_name": "planned.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/planned.md",
        "mime_type": "text/markdown",
        "size_bytes": 256,
        "checksum": "planned-checksum",
        "extracted_summary": "A source summary that is sufficient to generate a brief, outline, and slide plan.",
        "metadata": {},
    }

    with TestClient(app) as client:
        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief = brief_response.json()["brief"]

        outline_response = client.post(
            f"/projects/{project_id}/outline:generate",
            json={"brief_id": brief["id"]},
        )
        assert outline_response.status_code == 200
        outline = outline_response.json()["outline"]

        slide_plan_response = client.post(
            f"/projects/{project_id}/slide-plan:generate",
            json={"outline_id": outline["id"]},
        )
        assert slide_plan_response.status_code == 200
        slide_plan = slide_plan_response.json()["slide_plan"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

        status_response = client.get(f"/projects/{project_id}/status")
        assert status_response.status_code == 200
        status_payload = status_response.json()

        exports_response = client.get(f"/projects/{project_id}/exports")
        assert exports_response.status_code == 200
        exports_payload = exports_response.json()

    assert detail_payload["project"]["status"] == "planned"
    assert detail_payload["latest_brief"]["id"] == brief["id"]
    assert detail_payload["latest_outline"]["id"] == outline["id"]
    assert detail_payload["latest_slide_plan"]["id"] == slide_plan["id"]
    assert detail_payload["latest_artifact"] is None
    assert detail_payload["latest_export"] is None

    assert status_payload["project_status"] == "planned"
    assert status_payload["current_task"]["id"] == status_payload["recent_tasks"][-1]["id"]
    assert status_payload["current_task"]["task_type"] == "generate_slide_plan"
    assert status_payload["current_task"]["task_status"] == "succeeded"
    assert status_payload["current_task"]["result"]["slide_plan_id"] == slide_plan["id"]
    assert status_payload["current_task"]["result"]["outline_id"] == outline["id"]
    assert status_payload["current_task"]["result"]["brief_id"] == brief["id"]
    assert status_payload["recent_tasks"][-1]["task_type"] == "generate_slide_plan"
    assert status_payload["recent_tasks"][-1]["result"]["slide_plan_id"] == slide_plan["id"]
    assert status_payload["recent_tasks"][-1]["result"]["outline_id"] == outline["id"]
    assert status_payload["recent_tasks"][-1]["result"]["brief_id"] == brief["id"]
    assert exports_payload["exports"] == []


def test_template_override_can_regenerate_slide_plan_and_artifact(tmp_path: Path) -> None:
    database_path = tmp_path / "template-override-test.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    from app.main import create_app

    app = create_app()

    project_payload = {
        "name": "Template Override Demo",
        "description": "Verify preferred template and artifact template overrides are respected.",
        "source_mode": "file",
        "tags": ["phase6", "template"],
    }

    file_payload = {
        "file_name": "demo.md",
        "file_type": "markdown",
        "storage_path": "storage/uploads/demo.md",
        "mime_type": "text/markdown",
        "size_bytes": 128,
        "checksum": "demo-checksum",
        "extracted_summary": "A compact source summary for template override tests.",
        "metadata": {},
    }

    with TestClient(app) as client:
        templates_response = client.get("/projects/templates")
        assert templates_response.status_code == 200
        templates = templates_response.json()["templates"]
        assert len(templates) >= 2
        preferred_template_id = templates[0]["template_id"]
        rerender_template_id = templates[1]["template_id"]

        create_response = client.post("/projects", json=project_payload)
        assert create_response.status_code == 201
        project_id = create_response.json()["project"]["id"]

        register_response = client.post(f"/projects/{project_id}/files", json=file_payload)
        assert register_response.status_code == 201

        brief_response = client.post(f"/projects/{project_id}/brief:generate", json={"force_regenerate": True})
        assert brief_response.status_code == 200
        brief_id = brief_response.json()["brief"]["id"]

        outline_response = client.post(f"/projects/{project_id}/outline:generate", json={"brief_id": brief_id})
        assert outline_response.status_code == 200
        outline_id = outline_response.json()["outline"]["id"]

        slide_plan_response = client.post(
            f"/projects/{project_id}/slide-plan:generate",
            json={
                "outline_id": outline_id,
                "preferred_template_id": preferred_template_id,
                "force_regenerate": True,
            },
        )
        assert slide_plan_response.status_code == 200
        slide_plan = slide_plan_response.json()["slide_plan"]

        artifact_response = client.post(
            f"/projects/{project_id}/artifact:generate",
            json={
                "slide_plan_id": slide_plan["id"],
                "template_id": rerender_template_id,
            },
        )
        assert artifact_response.status_code == 200
        artifact = artifact_response.json()["artifact"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail_payload = detail_response.json()

    latest_slide_plan = detail_payload["latest_slide_plan"]
    latest_artifact = detail_payload["latest_artifact"]

    assert latest_slide_plan["id"] == slide_plan["id"]
    assert latest_slide_plan["metadata"]["preferred_template_id"] == preferred_template_id
    assert artifact["id"] == latest_artifact["id"]
    assert latest_artifact["metadata"]["template_id"] == rerender_template_id
    assert latest_artifact["metadata"]["template_name"]
    assert latest_artifact["render_status"] in {"succeeded", "partial"}