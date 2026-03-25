from __future__ import annotations

import base64
import json
import os
import sys
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pypdf import PdfReader

from app.services.sample_catalog import SampleCatalogService


REPO_ROOT = Path(__file__).resolve().parents[3]
CATEGORY_SET = {
    "学术论文答辩",
    "企业战略汇报",
    "产品发布会",
    "培训课件",
    "政府汇报",
    "技术分享",
}
EXPORT_FORMAT_BY_SAMPLE_ID = {
    "academic-thesis-defense": "pptx",
    "enterprise-strategy-review": "pptx",
    "product-launch-deck": "pdf",
    "training-courseware": "pptx",
    "government-reporting-pack": "pdf",
    "technical-sharing-session": "pptx",
}


def _load_samples() -> list[dict[str, object]]:
    return SampleCatalogService(REPO_ROOT).list_samples()


def _build_user_intent_payload(sample: dict[str, object], source_text: str) -> dict[str, object]:
    first_line = next((line.strip() for line in source_text.splitlines() if line.strip()), sample["summary"])
    return {
        "audience": sample["brief"]["audience"],
        "scenario": sample["category"],
        "purpose": sample["brief"]["goal"],
        "desired_page_count": 10,
        "style_preferences": [sample["brief"]["style"]],
        "emphasize_points": [sample["summary"], sample["project_name"]],
        "constraints": [first_line[:120]],
    }


def _create_sample_project(client: TestClient, sample: dict[str, object], source_text: str) -> tuple[str, dict[str, object]]:
    user_intent_payload = _build_user_intent_payload(sample, source_text)
    project_payload = {
        "name": sample["project_name"],
        "description": sample["summary"],
        "source_mode": sample["source_mode"],
        "tags": ["phase7", "sample-smoke", sample["sample_id"]],
    }

    create_response = client.post("/projects", json=project_payload)
    assert create_response.status_code == 201
    project_id = create_response.json()["project"]["id"]

    if sample["source_mode"] == "chat":
        brief_response = client.post(
            f"/projects/{project_id}/brief:generate",
            json={
                "force_regenerate": True,
                "user_intent_override": user_intent_payload,
            },
        )
        assert brief_response.status_code == 200
    else:
        upload_response = client.post(
            f"/projects/{project_id}/files:upload",
            json={
                "file_name": f"{sample['sample_id']}.md",
                "file_type": "markdown",
                "mime_type": "text/markdown",
                "content_base64": base64.b64encode(source_text.encode("utf-8")).decode("ascii"),
                "metadata": {"sample_id": sample["sample_id"], "sample_category": sample["category"]},
            },
        )
        assert upload_response.status_code == 201
        file_id = upload_response.json()["file"]["id"]

        parse_response = client.post(
            f"/projects/{project_id}/files:parse",
            json={
                "file_ids": [file_id],
                "rebuild_bundle": True,
                "user_intent": user_intent_payload,
            },
        )
        assert parse_response.status_code == 200
        assert parse_response.json()["source_bundle"] is not None

        brief_response = client.post(
            f"/projects/{project_id}/brief:generate",
            json={"force_regenerate": True, "user_intent_override": user_intent_payload},
        )
        assert brief_response.status_code == 200

    brief = brief_response.json()["brief"]
    assert brief["presentation_goal"]
    assert brief["target_audience"] == sample["brief"]["audience"]
    assert brief["recommended_page_count"] >= 8

    outline_response = client.post(
        f"/projects/{project_id}/outline:generate",
        json={"brief_id": brief["id"]},
    )
    assert outline_response.status_code == 200
    outline = outline_response.json()["outline"]
    assert outline["chapters"]
    assert outline["metadata"]["generation_mode"] == "methodology_engine_v1"

    slide_plan_response = client.post(
        f"/projects/{project_id}/slide-plan:generate",
        json={
            "outline_id": outline["id"],
            "preferred_template_id": sample["recommended_template_id"],
        },
    )
    assert slide_plan_response.status_code == 200
    slide_plan = slide_plan_response.json()["slide_plan"]

    assert slide_plan["page_count"] >= 3
    assert slide_plan["design_direction"] == sample["recommended_template_id"]
    assert len(slide_plan["slides"]) == slide_plan["page_count"]

    return project_id, {"brief": brief, "outline": outline, "slide_plan": slide_plan}


def test_sample_registry_covers_all_phase_7_3_categories() -> None:
    registry_path = REPO_ROOT / "storage/projects/sample-registry.json"

    assert registry_path.exists()

    payload = json.loads(registry_path.read_text(encoding="utf-8"))
    samples = payload["samples"]

    assert len(samples) >= 6

    categories = {item["category"] for item in samples}
    assert categories == CATEGORY_SET


def test_sample_registry_entries_keep_required_fields_and_unique_ids() -> None:
    samples = _load_samples()

    sample_ids = [item["sample_id"] for item in samples]
    assert len(sample_ids) == len(set(sample_ids))

    for item in samples:
        assert item["project_name"]
        assert item["summary"]
        assert item["source_mode"] in {"chat", "file", "hybrid"}
        assert item["recommended_template_id"]
        assert item["brief"]["audience"]
        assert item["brief"]["goal"]
        assert item["brief"]["style"]


def test_sample_registry_templates_exist_in_builtin_catalog() -> None:
    samples = _load_samples()

    builtin_template_ids = {path.stem for path in (REPO_ROOT / "templates/builtin").glob("*.json")}

    assert builtin_template_ids

    missing = {
        item["recommended_template_id"]
        for item in samples
        if item["recommended_template_id"] not in builtin_template_ids
    }
    assert missing == set()


def test_sample_registry_source_assets_exist_and_match_modes() -> None:
    samples = _load_samples()

    for item in samples:
        source_asset = item["source_asset"]
        assert isinstance(source_asset, dict)

        asset_kind = source_asset["kind"]
        asset_path = REPO_ROOT / source_asset["path"]

        assert asset_kind in {"inline-markdown", "prompt-markdown"}
        assert asset_path.exists()
        assert asset_path.suffix == ".md"
        assert asset_path.read_text(encoding="utf-8").strip()

        if item["source_mode"] == "file":
            assert asset_kind == "inline-markdown"
        if item["source_mode"] == "chat":
            assert asset_kind == "prompt-markdown"


def test_sample_catalog_service_returns_registry_entries_and_source_text() -> None:
    catalog = SampleCatalogService(REPO_ROOT)
    samples = catalog.list_samples()

    assert len(samples) >= 6

    first = catalog.get_sample(samples[0]["sample_id"])
    source_text = catalog.read_source_text(first)

    assert first["sample_id"] == samples[0]["sample_id"]
    assert source_text.strip()


@pytest.mark.parametrize("sample", _load_samples(), ids=lambda sample: sample["sample_id"])
def test_sample_registry_generation_smoke(sample: dict[str, object], tmp_path: Path) -> None:
    database_path = tmp_path / f"{sample['sample_id']}.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()
    catalog = SampleCatalogService(REPO_ROOT)
    source_text = catalog.read_source_text(sample)

    with TestClient(app) as client:
        project_id, generated = _create_sample_project(client, sample, source_text)
        brief = generated["brief"]
        outline = generated["outline"]
        slide_plan = generated["slide_plan"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail = detail_response.json()

    assert detail["project"]["status"] == "planned"
    assert detail["latest_brief"]["id"] == brief["id"]
    assert detail["latest_outline"]["id"] == outline["id"]
    assert detail["latest_slide_plan"]["id"] == slide_plan["id"]


@pytest.mark.parametrize(
    "sample_id,export_format",
    [
        ("academic-thesis-defense", "pptx"),
        ("product-launch-deck", "pdf"),
    ],
    ids=["academic-thesis-defense-pptx", "product-launch-deck-pdf"],
)
def test_sample_registry_render_export_smoke(sample_id: str, export_format: str, tmp_path: Path) -> None:
    database_path = tmp_path / f"{sample_id}-{export_format}.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()
    catalog = SampleCatalogService(REPO_ROOT)
    sample = catalog.get_sample(sample_id)
    source_text = catalog.read_source_text(sample)

    with TestClient(app) as client:
        project_id, generated = _create_sample_project(client, sample, source_text)
        slide_plan = generated["slide_plan"]

        render_response = client.post(
            f"/projects/{project_id}/artifact:generate",
            json={
                "slide_plan_id": slide_plan["id"],
                "template_id": sample["recommended_template_id"],
            },
        )
        assert render_response.status_code == 200
        artifact = render_response.json()["artifact"]
        assert artifact["metadata"]["template_id"] == sample["recommended_template_id"]
        assert artifact["render_status"] in {"succeeded", "partial"}
        assert artifact["metadata"]["generated_svg_files"]

        export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact["id"], "export_format": export_format},
        )
        assert export_response.status_code == 200
        export_payload = export_response.json()

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail = detail_response.json()

    export_job = export_payload["export_job"]
    export_path = Path(export_job["export_path"])
    assert export_path.exists()
    assert export_job["status"] == "succeeded"
    assert export_job["artifact_id"] == artifact["id"]
    assert export_job["run_id"].startswith("run-")
    assert export_job["metadata"]["file_count"] >= 1
    assert len(artifact["metadata"]["generated_svg_files"]) == export_job["metadata"]["file_count"]

    if export_format == "pptx":
        presentation = Presentation(export_path)
        assert export_path.suffix == ".pptx"
        assert export_job["metadata"]["export_kind"] == "pptx_from_svg_pages"
        assert len(presentation.slides) == export_job["metadata"]["file_count"]
    else:
        preview_pdf_path = Path(export_job["preview_pdf_path"])
        reader = PdfReader(str(preview_pdf_path))
        assert export_path.suffix == ".pdf"
        assert export_job["metadata"]["export_kind"] == "pdf_preview_from_svg_pages"
        assert preview_pdf_path == export_path
        assert len(reader.pages) == export_job["metadata"]["file_count"]

    assert detail["project"]["status"] == "exported"
    assert detail["latest_artifact"]["id"] == artifact["id"]
    assert detail["latest_export"]["id"] == export_job["id"]
    assert detail["latest_export"]["run_id"] == export_job["run_id"]


@pytest.mark.parametrize("sample", _load_samples(), ids=lambda sample: f"render-{sample['sample_id']}")
def test_sample_registry_render_smoke_all_samples(sample: dict[str, object], tmp_path: Path) -> None:
    database_path = tmp_path / f"render-{sample['sample_id']}.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()
    catalog = SampleCatalogService(REPO_ROOT)
    source_text = catalog.read_source_text(sample)

    with TestClient(app) as client:
        project_id, generated = _create_sample_project(client, sample, source_text)
        slide_plan = generated["slide_plan"]

        render_response = client.post(
            f"/projects/{project_id}/artifact:generate",
            json={"slide_plan_id": slide_plan["id"], "template_id": sample["recommended_template_id"]},
        )
        assert render_response.status_code == 200
        artifact = render_response.json()["artifact"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail = detail_response.json()

    generated_svg_files = artifact["metadata"]["generated_svg_files"]
    validation_summary = artifact["metadata"]["validation_summary"]
    finalization_summary = artifact["metadata"]["finalization_summary"]

    assert artifact["slide_plan_id"] == slide_plan["id"]
    assert artifact["metadata"]["template_id"] == sample["recommended_template_id"]
    assert artifact["render_status"] in {"succeeded", "partial"}
    assert generated_svg_files
    assert len(generated_svg_files) == slide_plan["page_count"]
    assert validation_summary["checked_file_count"] == slide_plan["page_count"]
    assert validation_summary["invalid_file_count"] == 0
    assert finalization_summary["page_count"] == slide_plan["page_count"]

    for svg_path in generated_svg_files:
        assert Path(svg_path).exists()

    assert detail["project"]["status"] == "finalized"
    assert detail["latest_slide_plan"]["id"] == slide_plan["id"]
    assert detail["latest_artifact"]["id"] == artifact["id"]
    assert detail["latest_artifact"]["metadata"]["template_id"] == sample["recommended_template_id"]
    assert detail["latest_export"] is None

@pytest.mark.parametrize("sample", _load_samples(), ids=lambda sample: f"export-{sample['sample_id']}")
def test_sample_registry_export_smoke_all_samples(sample: dict[str, object], tmp_path: Path) -> None:
    export_format = EXPORT_FORMAT_BY_SAMPLE_ID[sample["sample_id"]]
    database_path = tmp_path / f"export-{sample['sample_id']}.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()
    catalog = SampleCatalogService(REPO_ROOT)
    source_text = catalog.read_source_text(sample)

    with TestClient(app) as client:
        project_id, generated = _create_sample_project(client, sample, source_text)
        slide_plan = generated["slide_plan"]

        render_response = client.post(
            f"/projects/{project_id}/artifact:generate",
            json={"slide_plan_id": slide_plan["id"], "template_id": sample["recommended_template_id"]},
        )
        assert render_response.status_code == 200
        artifact = render_response.json()["artifact"]

        export_response = client.post(
            f"/projects/{project_id}/export",
            json={"artifact_id": artifact["id"], "export_format": export_format},
        )
        assert export_response.status_code == 200
        export_job = export_response.json()["export_job"]

        detail_response = client.get(f"/projects/{project_id}")
        assert detail_response.status_code == 200
        detail = detail_response.json()

    export_path = Path(export_job["export_path"])
    generated_svg_files = artifact["metadata"]["generated_svg_files"]

    assert artifact["metadata"]["template_id"] == sample["recommended_template_id"]
    assert export_job["status"] == "succeeded"
    assert export_job["artifact_id"] == artifact["id"]
    assert export_job["metadata"]["file_count"] == slide_plan["page_count"]
    assert len(generated_svg_files) == slide_plan["page_count"]
    assert export_path.exists()

    if export_format == "pptx":
        presentation = Presentation(export_path)
        preview_pdf_path_value = export_job.get("preview_pdf_path")
        assert export_path.suffix == ".pptx"
        assert export_job["metadata"]["export_kind"] == "pptx_from_svg_pages"
        assert len(presentation.slides) == slide_plan["page_count"]
        if preview_pdf_path_value:
            preview_pdf_path = Path(preview_pdf_path_value)
            assert preview_pdf_path.exists()
            assert len(PdfReader(str(preview_pdf_path)).pages) == slide_plan["page_count"]
    else:
        preview_pdf_path = Path(export_job["preview_pdf_path"])
        reader = PdfReader(str(export_path))
        assert export_path.suffix == ".pdf"
        assert preview_pdf_path == export_path
        assert export_job["metadata"]["export_kind"] == "pdf_preview_from_svg_pages"
        assert len(reader.pages) == slide_plan["page_count"]

    assert detail["project"]["status"] == "exported"
    assert detail["latest_artifact"]["id"] == artifact["id"]
    assert detail["latest_export"]["id"] == export_job["id"]
    assert detail["latest_export"]["artifact_id"] == artifact["id"]
    assert detail["latest_export"]["metadata"]["file_count"] == slide_plan["page_count"]


@pytest.mark.parametrize("sample", _load_samples(), ids=lambda sample: f"golden-{sample['sample_id']}")
def test_sample_registry_render_metadata_golden_all_samples(sample: dict[str, object], tmp_path: Path) -> None:
    database_path = tmp_path / f"golden-{sample['sample_id']}.db"

    os.environ["DATABASE_URL"] = f"sqlite:///{database_path.as_posix()}"
    os.environ["AUTO_CREATE_TABLES"] = "true"

    for module_name in ("app.main", "app.db.session"):
        sys.modules.pop(module_name, None)

    from app.main import create_app

    app = create_app()
    catalog = SampleCatalogService(REPO_ROOT)
    source_text = catalog.read_source_text(sample)

    with TestClient(app) as client:
        project_id, generated = _create_sample_project(client, sample, source_text)
        slide_plan = generated["slide_plan"]

        render_response = client.post(
            f"/projects/{project_id}/artifact:generate",
            json={"slide_plan_id": slide_plan["id"], "template_id": sample["recommended_template_id"]},
        )
        assert render_response.status_code == 200
        artifact = render_response.json()["artifact"]

    render_root = Path("storage") / "projects" / project_id / "render" / artifact["id"]
    design_spec_path = render_root / "design-spec.json"
    render_log_path = render_root / "render-log.json"
    svg_final_root = render_root / "svg_final"

    assert design_spec_path.exists()
    assert render_log_path.exists()
    assert svg_final_root.exists()

    design_spec_payload = json.loads(design_spec_path.read_text(encoding="utf-8"))
    render_log_payload = json.loads(render_log_path.read_text(encoding="utf-8"))
    finalized_svgs = list(svg_final_root.glob("slide-*.svg"))

    assert artifact["metadata"]["design_spec_path"].endswith("design-spec.json")
    assert design_spec_payload["project_id"] == project_id
    assert design_spec_payload["page_count"] == slide_plan["page_count"]
    assert design_spec_payload["template"]["template_id"] == sample["recommended_template_id"]
    assert render_log_payload["finalization_summary"]["page_count"] == slide_plan["page_count"]
    assert render_log_payload["validation_results"]
    assert len(render_log_payload["validation_results"]) == slide_plan["page_count"]
    assert all(item["is_valid"] for item in render_log_payload["validation_results"])
    assert all(item["finalizer_steps"] for item in render_log_payload["validation_results"])
    assert len(finalized_svgs) == slide_plan["page_count"]